# -*- coding: utf-8 -*-
"""
Created on Wed Jul  3 16:10:24 2024

@author: benny.ng
"""


from workalendar.america import Brazil
from pandas.tseries.offsets import CustomBusinessDay
import requests
import json
import pandas as pd
from pptx import Presentation
import comtypes.client
import os
import numpy as np
import locale
import yfinance as yahooFinance
from datetime import datetime, timedelta
from dateutil.relativedelta import relativedelta
from pptx.chart.data import CategoryChartData
from pptx.util import Pt, Cm
from copy import deepcopy
from pptx.table import Table, _Row, _Cell
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

locale.setlocale(locale.LC_ALL, 'pt_BR.UTF-8')


data_fim = datetime(2024, 6, 28)  # TODO MUDAR
folder = "JUNHO24"  # TODO MUDAR
username = "ekho.fo"
password = "EKH@fo2024"


portifolio_onshore = "1002"
data_ini_onshore = "23102023"

portifolio_offshore = "2000"
data_ini_offshore = "29112023"

client_name = "CARLOS EDUARDO DE ALMEIDA MOYSES_FexExplodido"


# Azul Claro, Laranja, Azul Escuro, Cinza, Verde, Preto, Branco
colors = ["1178ce", "ce6611", "002d74", "5b6670", "c2d500", "222b32", "c1c5c8"]


def pptx_to_pdf(pptx_folder, pdf_folder):
    # Inicia o PowerPoint
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1

    for filename in os.listdir(pptx_folder):
        if filename.endswith(".pptx"):
            file_path = os.path.join(pptx_folder, filename)
            pdf_path = os.path.join(
                pdf_folder, filename.replace(".pptx", ".pdf"))

            # Abre a apresentação
            presentation = powerpoint.Presentations.Open(
                file_path, WithWindow=False)

            # Salva como PDF
            presentation.SaveAs(pdf_path, 32)  # 32 representa o formato PDF
            presentation.Close()

    powerpoint.Quit()


def json_to_df(json, table):
    dict_table = json["tables"][table]

    df_table = pd.DataFrame.from_dict(dict_table, orient='index')
    # Converting strings with commas as decimal points to floats
    # Skipping the first column as it's categorical
    for col in df_table.columns[1:]:
        df_table[col] = df_table[col].apply(lambda x: float(x.replace(',', '.')) if x.replace(
            ',', '.').replace('-', '').replace('.', '').isdigit() else x)
    df_table.columns = df_table.iloc[0]
    df_table = df_table.drop(df_table.index[0])

    return df_table


def CMD_request(extract, portifolio, data_ini, data_fim, data_ini2, layout="1", extra="", benchmarks="CDI%2BIBOV%2BANBIMA_IMAB", mes_benchs="CDI%2Bpercent_CDI", ret2="ativo%2Bmes_atual%2Bano_atual%2B03m%2B06m%2B12m%2B24m%2B36m%2Bperiodo_compra%2Bdata_aplicacao%2Bpercent_SB", ret3="ativo%2Bpercent_SB%2Bmes_atual%2Bano_atual%2B12_m%2Bano_anterior%2Bperiodo_compra", username=username, password=password):
    url = "https://www.comdinheiro.com.br/Clientes/API/EndPoint001.php"
    querystring = {"code": "import_data"}
    payload = f"username={username}&password={password}&URL={extract}.php%3Fnome_portfolio%3D{portifolio}%26data_ini%3D{data_ini}%26data_fim%3D{data_fim}%26data_ini2%3D{data_ini2}{extra}%26layout%3D{layout}%26benchmarks%3D{benchmarks}%26mes_benchs%3D{mes_benchs}%26ret2%3D{ret2}%26ret3%3D{ret3}%26classe%3DMV%28class_ekho%29%26classe2%3DIF&format=json3"
    headers = {'Content-Type': 'application/x-www-form-urlencoded'}
    response = requests.request(
        "POST", url, data=payload, headers=headers, params=querystring)

    json_carteira = json.loads(response.text)

    return json_carteira


def delete_slide(presentation, slide_index):
    # Acesso direto ao objeto XML que contém os slides
    xml_slides = presentation.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[slide_index])  # Remove o slide do xml


def round_if_numeric(x, casas, string="", before=True):
    try:
        number = round(pd.to_numeric(x), casas)
        formated = locale.format_string(
            f"%.{casas}f", number, grouping=True) if casas != 0 else locale.format_string("%.0f", number, grouping=True)
        formated = f"{string}{formated}" if before == True else f"{formated}{string}"
        return formated
    except ValueError:
        return x


def first_workday(date):
    year = date.year
    month = date.month
    cal = Brazil()
    first_day = datetime(year, month, 1)

    while not cal.is_working_day(first_day):
        first_day += timedelta(days=1)
    return first_day


def previous_month_last_workday(date):
    year = date.year
    month = date.month
    cal = Brazil()
    last_day = datetime(year, month, 1)+timedelta(days=-1)

    while not cal.is_working_day(last_day, extra_holidays=[datetime(2024, 3, 29, 0, 0)]):
        last_day += timedelta(days=-1)
    return last_day


def accum_return(series, start_date, end_date):
    # start_date = "2024-01-01"
    # end_date = "2024-04-01"
    # series = df_benchmarks["GLAD"]
    name = series.name
    series.name = "Retorno_Diario"
    df = pd.DataFrame(series)
    selected_data = df[(df.index >= start_date) &
                       (df.index <= end_date)].copy()
    selected_data.iloc[0, 0] = 0
    selected_data['Fator'] = 1 + selected_data['Retorno_Diario']
    selected_data['Retorno_Acumulado'] = selected_data['Fator'].cumprod()
    selected_data['Retorno_Acumulado'] -= 1
    selected_data['Retorno_Acumulado'] = selected_data['Retorno_Acumulado']*100

    final = selected_data['Retorno_Acumulado']
    final.name = name

    return final


def convert_to_float(val):
    try:
        # Attempt to replace commas and convert to float
        return float(str(val).replace(',', '.'))
    except ValueError:
        # Return the value unchanged if conversion fails
        return val


def getList(dict):
    list = []
    for key in dict.keys():
        list.append(key)

    return list


meses = {
    1: "Janeiro",
    2: "Fevereiro",
    3: "Março",
    4: "Abril",
    5: "Maio",
    6: "Junho",
    7: "Julho",
    8: "Agosto",
    9: "Setembro",
    10: "Outubro",
    11: "Novembro",
    12: "Dezembro",
}

date_title = f"{meses[data_fim.month].upper()} | {data_fim.year}"
data_fim_str = data_fim.strftime("%d%m%Y")
classe_benchmark = pd.read_excel(
    open('Suporte_Relatorio.xlsx', 'rb'), sheet_name='classes')
carteira_ekho_onshore = pd.read_excel(
    open('Suporte_Relatorio.xlsx', 'rb'), sheet_name='carteira_ekho_onshore')
carteira_ekho_offshore = pd.read_excel(
    open('Suporte_Relatorio.xlsx', 'rb'), sheet_name='carteira_ekho_offshore')
classe_benchmark["class_region"] = classe_benchmark["class"] + \
    "_"+classe_benchmark["region"]
clientes_on_off = pd.read_excel(
    open('Suporte_Relatorio.xlsx', 'rb'), sheet_name='clients_on_off')
clientes_onshore = pd.read_excel(
    open('Suporte_Relatorio.xlsx', 'rb'), sheet_name='clients_onshore')

first_workday_str = first_workday(data_fim).strftime("%d%m%Y")
previous_month_last_workday_str = previous_month_last_workday(data_fim).strftime(
    "%d%m%Y")

# Benchmarks
ACWI = yahooFinance.Ticker("ACWI").history(
    period="max")["Close"]  # MSCI World ACWI
ACWI = ACWI.tz_localize(None)
df_ACWI = pd.DataFrame(ACWI)
df_ACWI["ACWI"] = df_ACWI['Close'].pct_change()

GLAG = yahooFinance.Ticker("GLAG.L").history(period="max")[
    "Close"]  # Bloomberg Global Aggregate Bonds Index
GLAG = GLAG.tz_localize(None)
df_GLAG = pd.DataFrame(GLAG)
df_GLAG["GLAG"] = df_GLAG['Close'].pct_change()


SPX = yahooFinance.Ticker("%5ESPX").history(period="max")[
    "Close"]  # S&P 500
SPX = SPX.tz_localize(None)
df_SPX = pd.DataFrame(SPX)
df_SPX["SPX"] = df_SPX['Close'].pct_change()

data_fim_benchmark = datetime.now().strftime("%d%m%Y")
url = "https://www.comdinheiro.com.br/Clientes/API/EndPoint001.php"
querystring = {"code": "import_data"}
payload = f"username={username}&password={password}&URL=HistoricoCotacao002.php%3F%26x%3DANBIMA_IMAB%2BANBIMA_IRFM%2BANBIMA_IHFA%2BCDI%2BIFIX%2BUS_EFFR%2BIBOV%26data_ini%3D05122008%26data_fim%3D{data_fim_benchmark}%26pagina%3D1%26d%3DMOEDA_ORIGINAL%26g%3D1%26m%3D0%26info_desejada%3Dretorno%26retorno%3Ddiscreto%26tipo_data%3Ddu_br%26tipo_ajuste%3Dtodosajustes%26num_casas%3D15%26enviar_email%3D0%26ordem_legenda%3D1%26cabecalho_excel%3Dmodo1%26classes_ativos%3Dz1ci99jj7473%26ordem_data%3D0%26rent_acum%3Drent_acum%26minY%3D%26maxY%3D%26deltaY%3D%26preco_nd_ant%3D0%26base_num_indice%3D100%26flag_num_indice%3D0%26eixo_x%3DData%26startX%3D0%26max_list_size%3D20%26line_width%3D2%26titulo_grafico%3D%26legenda_eixoy%3D%26tipo_grafico%3Dline%26script%3D%26tooltip%3Dunica&format=json3"
headers = {'Content-Type': 'application/x-www-form-urlencoded'}
response = requests.request(
    "POST", url, data=payload, headers=headers, params=querystring)
json_benchmarks = json.loads(response.text)
df_benchmarks = json_to_df(json_benchmarks, "tab1")
df_benchmarks['Data'] = pd.to_datetime(df_benchmarks['Data'], dayfirst=True)
df_benchmarks.set_index('Data', inplace=True)

df_benchmarks = pd.concat(
    [df_benchmarks, df_ACWI["ACWI"], df_GLAG["GLAG"], df_SPX["SPX"]], axis=1)
df_benchmarks = df_benchmarks.replace('nd', np.nan)
for coluna in df_benchmarks.columns:
    df_benchmarks[coluna] = df_benchmarks[coluna].apply(
        lambda x: float(x.replace(",", ".")) if isinstance(x, str) else x)


def benchmark_diff(ativo, benchmark, metric):
    if ativo == "nd":
        return"-"
    else:
        if metric == "%":
            return ((ativo/benchmark)*100)
        else:
            return (ativo-benchmark)


def add_row(table: Table) -> _Row:
    new_row = deepcopy(table._tbl.tr_lst[-1])
    # duplicating last row of the table as a new row to be added

    for tc in new_row.tc_lst:
        cell = _Cell(tc, new_row.tc_lst)
        cell.text = ''  # defaulting cell contents to empty text

    table._tbl.append(new_row)


def remove_row(table, row):
    tbl = table._tbl
    tr = row._tr
    tbl.remove(tr)


def fill_asset_table(start_slide, dict_df, max_rows_len=18):
    slide_counter = start_slide
    table_idx_list = []
    number_of_columns = 11

    for index in range(len(dict_df)):
        print(f"index: {index}; class: {getList(dict_df)[index]}")
        previous_table_rows_len = 0
        next_table = False

        while not next_table:
            table_rows = 20

            slide = prs.slides[slide_counter]  # Acessa o slide
            table_idx_list.append(index)  # Adiciona tabela percorrida na lista

            rentabilidade_ativo_table = [
                shape for shape in slide.shapes if shape.name == f"rentabilidade_ativo_{index}"][0].table  # Seleciona a tabela pptx

            df = dict_df[getList(dict_df)[index]]
            # Seleciona o DataFrame e corta caso já esteja em uso
            df = df.iloc[previous_table_rows_len:]

            if previous_table_rows_len != 0:
                table_rows -= 1
                # Remove a primeira linha caso a tabela esteja na metade
                remove_row(rentabilidade_ativo_table,
                           rentabilidade_ativo_table.rows[0])
            else:
                max_rows_len = max_rows_len - \
                    1 if len(table_idx_list)//2 == 1 else max_rows_len+1

            table_rows_len = max_rows_len if max_rows_len < len(
                df) else len(df)  # Determina numero de linhas na tabela
            next_table = True if table_rows_len == len(
                df) else False  # Verifica se pode ir para prox tabela
            # Verifica se é possível adicionar mais uma tabela
            max_rows_len = max_rows_len - \
                len(df) if max_rows_len-len(df) >= 3 else 0

            for i in range(table_rows_len):
                for j in range(number_of_columns):
                    rentabilidade_ativo_table.cell(
                        i, j+1).text_frame.paragraphs[0].runs[0].text = df.iloc[i, j]

            for row_idx in range(table_rows_len, table_rows+1):
                remove_row(rentabilidade_ativo_table,
                           rentabilidade_ativo_table.rows[table_rows_len])

            if max_rows_len == 0:
                # print("#######")
                for table_idx in range(8):  # remove as tabelas nao usadas
                    if table_idx not in table_idx_list:
                        shape = [shape for shape in slide.shapes if shape.name ==
                                 f"rentabilidade_ativo_{table_idx}"][0]
                        sp = shape._element
                        sp.getparent().remove(sp)

                for idx, elem in enumerate(table_idx_list):
                    shape = [shape for shape in slide.shapes if shape.name ==
                             f"rentabilidade_ativo_{elem}"][0]
                    shape.left = 1484154
                    if idx == 0:
                        shape.top = 932400
                    else:
                        prev_shape = [
                            shape for shape in slide.shapes if shape.name == f"rentabilidade_ativo_{elem-1}"][0]
                        shape.top = prev_shape.top + \
                            len(prev_shape.table.rows)*Cm(0.79) + Cm(0.25)

                previous_table_rows_len = previous_table_rows_len + table_rows_len
                max_rows_len = 18
                slide_counter += 1
                table_idx_list = []

    for table_idx in range(8):  # limpa as tabelas nao usadas
        if table_idx not in table_idx_list:
            # print(table_idx)
            slide = prs.slides[slide_counter]
            shape = [shape for shape in slide.shapes if shape.name ==
                     f"rentabilidade_ativo_{table_idx}"][0]
            sp = shape._element
            sp.getparent().remove(sp)
        if len(slide.shapes) == 4:
            slide_counter -= 1

    for idx, elem in enumerate(table_idx_list):
        # print(f"idx: {idx}")
        # print(f"elem: {elem}")
        shape = [shape for shape in slide.shapes if shape.name ==
                 f"rentabilidade_ativo_{elem}"][0]
        shape.left = 1484154
        if idx == 0:
            shape.top = 932400
        else:
            prev_shape = [shape for shape in slide.shapes if shape.name ==
                          f"rentabilidade_ativo_{elem-1}"][0]
            shape.top = prev_shape.top + \
                len(prev_shape.table.rows)*Cm(0.79) + Cm(0.25)
    return slide_counter


def cmd_historico_acum(username, password, ticker, data_ini, data_fim, num_casas=5):
    url = "https://www.comdinheiro.com.br/Clientes/API/EndPoint001.php"
    querystring = {"code": "import_data"}
    payload = f"username={username}&password={password}&URL=HistoricoCotacao002.php%3F%26x%3D{ticker}%26data_ini%3D{data_ini}%26data_fim%3D{data_fim}%26g%3D1%26m%3D0%26info_desejada%3Dretorno_acum%26num_casas%3D{num_casas}&format=json3"
    headers = {'Content-Type': 'application/x-www-form-urlencoded'}
    response = requests.request(
        "POST", url, data=payload, headers=headers, params=querystring)
    json_historico = json.loads(response.text)
    df_historico = json_to_df(json_historico, "tab1")
    df_historico["Data"] = pd.to_datetime(df_historico["Data"], dayfirst=True)
    max_date = df_historico["Data"].max()
    valor = df_historico[df_historico["Data"] == max_date][ticker][0]
    valor = 0 if "E-" in str(valor) else valor

    return valor


ipca = {}
ipca["mes"] = cmd_historico_acum(username, password, ticker="IBGE_IPCA15", data_ini=previous_month_last_workday(
    data_fim).strftime("%d%m%Y"), data_fim=data_fim_str)
ipca["ano"] = cmd_historico_acum(username, password, ticker="IBGE_IPCA15", data_ini=previous_month_last_workday(
    datetime(data_fim.year, 1, 1)).strftime("%d%m%Y"), data_fim=data_fim_str)
ipca["12m"] = cmd_historico_acum(username, password, ticker="IBGE_IPCA15", data_ini=(
    data_fim+relativedelta(months=-12)).strftime("%d%m%Y"), data_fim=data_fim_str)
ipca["24m"] = cmd_historico_acum(username, password, ticker="IBGE_IPCA15", data_ini=(
    data_fim+relativedelta(months=-24)).strftime("%d%m%Y"), data_fim=data_fim_str)


cpi = {}
cpi["mes"] = cmd_historico_acum(username, password, ticker="US:BLS_CPI_forecast",
                                data_ini=previous_month_last_workday(data_fim).strftime("%d%m%Y"), data_fim=data_fim_str)
cpi["ano"] = cmd_historico_acum(username, password, ticker="US:BLS_CPI_forecast", data_ini=previous_month_last_workday(
    datetime(data_fim.year, 1, 1)).strftime("%d%m%Y"), data_fim=data_fim_str)
cpi["12m"] = cmd_historico_acum(username, password, ticker="US:BLS_CPI_forecast", data_ini=(
    data_fim+relativedelta(months=-12)).strftime("%d%m%Y"), data_fim=data_fim_str)
cpi["24m"] = cmd_historico_acum(username, password, ticker="US:BLS_CPI_forecast", data_ini=(
    data_fim+relativedelta(months=-24)).strftime("%d%m%Y"), data_fim=data_fim_str)


def tax_diff(taxa1, taxa2):
    is_string = False if not(
        type(taxa1) == str or type(taxa2) == str) else True
    taxa1 = taxa1/100 if not (is_string) else 0
    taxa2 = taxa2/100 if not (is_string) else 0
    diff = 100*(((1+taxa1)/(1+taxa2))-1) if not (is_string) else "-"
    return diff


def get_cmd_data(portifolio, region="onshore", data_ini="02012000"):
    # PL

    url = "https://www.comdinheiro.com.br/Clientes/API/EndPoint001.php"
    querystring = {"code": "import_data"}
    payload = f"username={username}&password={password}&URL=RelatorioGerencialCarteiras009.php%3F%26nome_portfolio%3D{portifolio}%26data_analise%3D{data_fim_str}%26data_ini%3D{data_ini}&format=json3"
    headers = {'Content-Type': 'application/x-www-form-urlencoded'}
    response = requests.request(
        "POST", url, data=payload, headers=headers, params=querystring)
    json_pl = json.loads(response.text)
    df_PL = json_to_df(json_pl, "tab0")
    df_PL = df_PL[df_PL.PL != "nd"]
    df_PL["Data"] = pd.to_datetime(df_PL["Data"], dayfirst=True)
    df_PL = df_PL.set_index("Data")

    payload = f"username={username}&password={password}&URL=HistoricoCotacao002.php%3F%26x%3DExt_{portifolio}_Cota%26data_ini%3D{data_ini}%26data_fim%3D{data_fim_str}%26info_desejada%3Dretorno_acum&format=json3"
    response = requests.request(
        "POST", url, data=payload, headers=headers, params=querystring)
    json_cota = json.loads(response.text)
    df_cota = json_to_df(json_cota, "tab1")
    df_cota["Data"] = pd.to_datetime(df_cota["Data"], dayfirst=True)
    df_cota = df_cota.set_index("Data")
    df_cota.iloc[0] = 0
    df_cota.rename(columns={df_cota.columns[0]: "Cota"}, inplace=True)

    df_PL = pd.concat([df_cota, df_PL["PL"]], axis=1)

    data_inicio_carteira = pd.to_datetime(df_PL.index[0], dayfirst=True)
    data_inicio_carteira_str = data_inicio_carteira.strftime("%d%m%Y")

    json_e21 = CMD_request("ExtratoCarteira021", portifolio, data_inicio_carteira_str,
                           data_fim_str, previous_month_last_workday_str, layout="2", extra="%26mes_acumulado%3D1%26mes_formato%3DV", mes_benchs="percent_CDI")
    df_CarteiraBenchmark = json_to_df(json_e21, "tab0")
    df_FluxoResgate = json_to_df(json_e21, "tab1")
    df_ResumoPorConta = json_to_df(json_e21, "tab8")

    data_12m = data_fim + relativedelta(months=-12)
    is_more_than_12m = True if data_12m > data_inicio_carteira else False

    data_24m = data_fim + relativedelta(months=-24)
    is_more_than_24m = True if data_24m > data_inicio_carteira else False

    if not(is_more_than_12m):
        df_CarteiraBenchmark.iloc[0, 3] = df_CarteiraBenchmark.iloc[0, 5]
        df_CarteiraBenchmark.iloc[1, 3] = "-"

    if not(is_more_than_24m):
        df_CarteiraBenchmark.iloc[0, 4] = df_CarteiraBenchmark.iloc[0, 5]
        df_CarteiraBenchmark.iloc[1, 4] = "-"

    if region == "offshore":
        url = "https://www.comdinheiro.com.br/Clientes/API/EndPoint001.php"

        querystring = {"code": "import_data"}

        payload = f"username={username}&password={password}&URL=ExtratoCarteira021.php%3F%26nome_portfolio%3D{portifolio}%26mes_acumulado%3D1%26data_ini%3D{data_inicio_carteira_str}%26data_fim%3D{data_fim_str}%26mes_benchs%3Dpercent%28US_EFFR%29&format=json3"
        headers = {'Content-Type': 'application/x-www-form-urlencoded'}

        response = requests.request(
            "POST", url, data=payload, headers=headers, params=querystring)

        json_e21 = json.loads(response.text)

    else:
        json_e21 = json_e21

    df_RetornosNominais = json_to_df(json_e21, "tab2")
    df_RetornosNominais = df_RetornosNominais[df_RetornosNominais.Ano != "nd"]
    # df_RetornosNominais = df_RetornosNominais[np.arange(
    #     len(df_RetornosNominais)) % 3 != 1]
    df_RetornosNominais = df_RetornosNominais.reset_index(drop=True)
    for index, row in df_RetornosNominais.iterrows():
        # Check if the row is for "Carteira"
        if row['Ativo'] == 'Carteira':
            # Iterate over each column (excluding the 'Ativo' and 'Jan' columns)
            for col in df_RetornosNominais.columns[2:]:
                # Convert to numeric, errors='coerce' will convert non-numeric values to NaN
                value = pd.to_numeric(
                    df_RetornosNominais.at[index, col], errors='coerce')
                # If the value is a number (not NaN)
                if not pd.isnull(value):
                    # Multiply the numeric value by 100
                    df_RetornosNominais.at[index, col] = value * 100
    for index, row in df_RetornosNominais.iterrows():
        # Check if the row is for "Carteira"
        if row['Ativo'] == '%CDI':
            # Iterate over each column (excluding the 'Ativo' and 'Jan' columns)
            for col in df_RetornosNominais.columns[2:]:
                # Convert to numeric, errors='coerce' will convert non-numeric values to NaN
                value = pd.to_numeric(
                    df_RetornosNominais.at[index, col], errors='coerce')
                # If the value is a number (not NaN)
                if not pd.isnull(value):
                    # Multiply the numeric value by 100
                    df_RetornosNominais.at[index, col] = value * 100

    for index, row in df_RetornosNominais.iterrows():
        if row['Ativo'] == '%US_EFFR':
            # Iterate over each column (excluding the 'Ativo' and 'Jan' columns)
            for col in df_RetornosNominais.columns[2:]:
                # Convert to numeric, errors='coerce' will convert non-numeric values to NaN
                value = pd.to_numeric(
                    df_RetornosNominais.at[index, col], errors='coerce')
                # If the value is a number (not NaN)
                if not pd.isnull(value):
                    # Multiply the numeric value by 100
                    df_RetornosNominais.at[index, col] = value * 100

    json_e22 = CMD_request("ExtratoCarteira022", portifolio, data_inicio_carteira_str,
                           data_fim_str, first_workday_str, "1")
    df_RetornoNominalAtivo = json_to_df(json_e22, "tab4")
    for i in range(len(df_RetornoNominalAtivo)):
        for j in range(len(df_RetornoNominalAtivo.iloc[i])):
            if type(df_RetornoNominalAtivo.iloc[i, j]) == str and "E-" in df_RetornoNominalAtivo.iloc[i, j]:
                df_RetornoNominalAtivo.iloc[i, j] = 0
            else:
                pass

    json_e22l6 = CMD_request(
        "ExtratoCarteira022", portifolio, data_inicio_carteira_str, data_fim_str, "mes_atual", "6", "%26data_ini3%3Dano_atual")
    df_PerfAttr = json_to_df(json_e22l6, "tab2")

    df_PerfAttr.iloc[:, :] = df_PerfAttr.iloc[:, :].applymap(convert_to_float)
    
    
    json_e22l6_ano = CMD_request(
        "ExtratoCarteira022", portifolio, data_inicio_carteira_str, data_fim_str, "ano_atual", "6", "%26data_ini3%3Dano_atual")
    df_PerfAttr_ano = json_to_df(json_e22l6_ano, "tab2")

    df_PerfAttr_ano.iloc[:, :] = df_PerfAttr_ano.iloc[:, :].applymap(convert_to_float)

    json_e23 = CMD_request("ExtratoCarteira023", portifolio,
                           data_inicio_carteira_str, data_fim_str, first_workday_str, ret3="ativo%2Bmes_atual%2Bano_atual%2B6m%2B12m%2B24m%2B36m%2Bpercent_SB%2Bsaldo_bruto")
    df_ClasseDeAtivos = json_to_df(json_e23, "tab1")

    json_e7 = CMD_request("ExtratoCarteira007", portifolio, data_inicio_carteira_str,
                          data_fim_str, first_workday_str, "1")
    df_Volatilidade = json_to_df(json_e7, "tab3")

    ipca_desde_ini = cmd_historico_acum(
        username, password, ticker="IBGE_IPCA15", data_ini=data_inicio_carteira_str, data_fim=data_fim_str)

    cpi_desde_ini = cmd_historico_acum(
        username, password, ticker="US:BLS_CPI_forecast", data_ini=data_inicio_carteira_str, data_fim=data_fim_str)
    

    cmd_data = {"data_inicio_carteira": data_inicio_carteira,
                "df_PL": df_PL,
                "df_CarteiraBenchmark": df_CarteiraBenchmark,
                "df_FluxoResgate": df_FluxoResgate,
                "df_RetornosNominais": df_RetornosNominais,
                "df_ResumoPorConta": df_ResumoPorConta,
                "df_RetornoNominalAtivo": df_RetornoNominalAtivo,
                "df_PerfAttr": df_PerfAttr,
                "df_PerfAttr_ano": df_PerfAttr_ano,
                "df_ClasseDeAtivos": df_ClasseDeAtivos,
                "df_Volatilidade": df_Volatilidade,
                "ipca_desde_ini": ipca_desde_ini,
                "cpi_desde_ini": cpi_desde_ini}

    return cmd_data


def perf_attr_fill(order_list, df_PerfAttr, slide):

    slide_width = Cm(33.87)
    df_PerfAttr = df_PerfAttr[df_PerfAttr["Classe"] != "Carteira"]
    carteira = df_PerfAttr["Bps no Mês"].sum()

    if "Caixa" in order_list:
        order_list.remove("Caixa")

    df_PerfAttr.loc[:, 'Classe_cat'] = pd.Categorical(
        df_PerfAttr['Classe'],
        categories=order_list,
        ordered=True)
    df_PerfAttr = df_PerfAttr.sort_values('Classe_cat').reset_index(drop=True)

    max_abs_value = 0
    sum_value = 0
    for i in range(len(df_PerfAttr)):
        sum_value = sum_value + df_PerfAttr.iloc[i, 1]
        if abs(sum_value) > max_abs_value:
            max_abs_value = abs(sum_value)
    df_PerfAttr["Ajustado"] = df_PerfAttr["Bps no Mês"]/max_abs_value
    carteira_ajustado = carteira/max_abs_value
    for i in range(len(df_PerfAttr)):
        shape = [shape for shape in slide.shapes if shape.name ==
                 f"perf_{i}"][0]
        shape.height = int(shape.height*abs(df_PerfAttr["Ajustado"][i]))
    for i in range(len(df_PerfAttr), 7):
        shape = [shape for shape in slide.shapes if shape.name ==
                 f"perf_{i}"][0]
        class_ekho = [
            shape for shape in slide.shapes if shape.name == f"class_{i}"][0]
        label = [shape for shape in slide.shapes if shape.name ==
                 f"label_{i}"][0]
        class_ekho.text_frame.paragraphs[0].runs[0].text = ""
        label.text_frame.paragraphs[0].runs[0].text = ""
        shape.width = 0
        class_ekho.width = 0
        label.width = 0

    shape = [shape for shape in slide.shapes if shape.name == "perf_total"][0]
    default_height = shape.top+shape.height
    shape.height = int(shape.height*abs(carteira_ajustado))
    shape.top = shape.top + default_height - (shape.top + shape.height)

    previous_shape = [
        shape for shape in slide.shapes if shape.name == "perf_total"][0]
    previous_class_value = carteira_ajustado
    i = len(df_PerfAttr)-1
    shape = [shape for shape in slide.shapes if shape.name == f"perf_{i}"][0]
    if (df_PerfAttr["Ajustado"][i] < 0 and previous_class_value > 0):
        shape.top = previous_shape.top - shape.height

    elif (df_PerfAttr["Ajustado"][i] > 0 and previous_class_value > 0):
        shape.top = previous_shape.top

    elif (df_PerfAttr["Ajustado"][i] < 0 and previous_class_value < 0):
        shape.top = previous_shape.top + (previous_shape.height - shape.height)

    elif (df_PerfAttr["Ajustado"][i] > 0 and previous_class_value < 0):
        shape.top = previous_shape.top + previous_shape.height

    for i in range(len(df_PerfAttr)-2, -1, -1):
        previous_shape = [
            shape for shape in slide.shapes if shape.name == f"perf_{i+1}"][0]
        previous_class_value = df_PerfAttr["Ajustado"][i+1]

        shape = [shape for shape in slide.shapes if shape.name ==
                 f"perf_{i}"][0]
        valor = df_PerfAttr["Ajustado"][i]

        if (valor < 0 and previous_class_value > 0):
            shape.top = previous_shape.top + \
                (previous_shape.height - shape.height)

        elif (valor > 0 and previous_class_value > 0):
            shape.top = previous_shape.top + previous_shape.height

        elif (valor < 0 and previous_class_value < 0):
            shape.top = previous_shape.top - shape.height

        elif (valor > 0 and previous_class_value < 0):
            shape.top = previous_shape.top

    total_bars = len(df_PerfAttr)+1
    total_width = int(total_bars*shape.width + (total_bars-1)*(Cm(1)))
    margin = int((slide_width - total_width)/2)
    shape = [shape for shape in slide.shapes if shape.name == f"perf_{0}"][0]
    class_ekho = [
        shape for shape in slide.shapes if shape.name == f"class_{0}"][0]
    class_ekho.text_frame.paragraphs[0].runs[0].text = df_PerfAttr["Classe"][0]
    label = [shape for shape in slide.shapes if shape.name == f"label_{0}"][0]
    label.text_frame.paragraphs[0].runs[0].text = str(
        int(np.round(df_PerfAttr["Bps no Mês"][0])))
    shape.left = margin
    class_ekho.left = margin
    label.left = margin
    next_margin = shape.left + shape.width + Cm(1)
    for i in range(1, len(df_PerfAttr)):
        shape = [shape for shape in slide.shapes if shape.name ==
                 f"perf_{i}"][0]
        class_ekho = [
            shape for shape in slide.shapes if shape.name == f"class_{i}"][0]
        class_ekho.text_frame.paragraphs[0].runs[0].text = df_PerfAttr["Classe"][i]
        label = [shape for shape in slide.shapes if shape.name ==
                 f"label_{i}"][0]
        label.text_frame.paragraphs[0].runs[0].text = str(
            int(np.round(df_PerfAttr["Bps no Mês"][i])))
        shape.left = next_margin
        class_ekho.left = next_margin
        label.left = next_margin
        next_margin = next_margin + shape.width + Cm(1)

    shape = [shape for shape in slide.shapes if shape.name == "perf_total"][0]
    class_ekho = [
        shape for shape in slide.shapes if shape.name == "class_total"][0]
    class_ekho.text_frame.paragraphs[0].runs[0].text = "Total da Carteira"
    label = [shape for shape in slide.shapes if shape.name == "label_total"][0]
    label.text_frame.paragraphs[0].runs[0].text = str(
        int(np.round(df_PerfAttr["Bps no Mês"].sum())))
    shape.left = next_margin
    class_ekho.left = next_margin
    label.left = next_margin

    for i in range(len(df_PerfAttr)):
        shape = [shape for shape in slide.shapes if shape.name ==
                 f"perf_{i}"][0]
        label = [shape for shape in slide.shapes if shape.name ==
                 f"label_{i}"][0]
        if df_PerfAttr["Bps no Mês"][i] > 0:
            label.top = shape.top-Cm(0.6)
        else:
            label.top = shape.top + shape.height

    shape = [shape for shape in slide.shapes if shape.name == "perf_total"][0]
    label = [shape for shape in slide.shapes if shape.name == "label_total"][0]
    if df_PerfAttr["Bps no Mês"].sum() > 0:
        label.top = shape.top-Cm(0.6)
    else:
        label.top = shape.top + shape.height
    
    shape_height = []
    for i in range(len(df_PerfAttr)):
        shape = [shape for shape in slide.shapes if shape.name ==
                 f"perf_{i}"][0]
        shape_height.append(shape.height + shape.top)
    if np.max(shape_height) + Cm(0.6) > Cm(9.37):
        max_diff = np.max(shape_height) + Cm(0.6) - Cm(9.37)
        for i in range(len(df_PerfAttr)):
            shape = [shape for shape in slide.shapes if shape.name ==
                     f"perf_{i}"][0]
            label = [shape for shape in slide.shapes if shape.name == f"label_{i}"][0]
            shape.top = shape.top-max_diff 
            label.top = label.top-max_diff 
        shape = [shape for shape in slide.shapes if shape.name == "perf_total"][0]
        label = [shape for shape in slide.shapes if shape.name == "label_total"][0]
        shape.top = shape.top-max_diff 
        label.top = label.top-max_diff 
        



#TODO TESTE
cmd_onshore = get_cmd_data(
    portifolio_onshore, region="onshore", data_ini=data_ini_onshore)
cmd_offshore = get_cmd_data(
    portifolio_offshore, region="offshore", data_ini=data_ini_offshore)

ipca["desde_ini"] = cmd_onshore["ipca_desde_ini"]
cpi["desde_ini"] = cmd_offshore["cpi_desde_ini"]

final_onshore = {}
final_offshore = {}
final_onshore["start_date"] = cmd_onshore["df_PL"].index[0]
final_onshore["end_date"] = cmd_onshore["df_PL"].index[-1]
final_offshore["start_date"] = cmd_offshore["df_PL"].index[0]
final_offshore["end_date"] = cmd_offshore["df_PL"].index[-1]
######################## TRATAMENTO DOS DADOS #########################
# SLIDE 1
final_onshore["df_final_CarteiraBenchmark"] = cmd_onshore["df_CarteiraBenchmark"].copy(
)
final_onshore["df_final_CarteiraBenchmark"].replace(
    'nd', '-', inplace=True)
final_onshore["df_final_CarteiraBenchmark"].iloc[4,
                                                  :] = final_onshore["df_final_CarteiraBenchmark"].iloc[3, :]
final_onshore["df_final_CarteiraBenchmark"].iloc[3,
                                                  :] = final_onshore["df_final_CarteiraBenchmark"].iloc[2, :]
final_onshore["df_final_CarteiraBenchmark"].iloc[2,
                                                  0] = "Retorno Real"
final_onshore["df_final_CarteiraBenchmark"].iloc[2, 1] = tax_diff(
    final_onshore["df_final_CarteiraBenchmark"].iloc[1, 1], ipca["mes"])  # MÊS
final_onshore["df_final_CarteiraBenchmark"].iloc[2, 2] = tax_diff(
    final_onshore["df_final_CarteiraBenchmark"].iloc[1, 2], ipca["ano"])  # ANO
final_onshore["df_final_CarteiraBenchmark"].iloc[2, 3] = tax_diff(
    final_onshore["df_final_CarteiraBenchmark"].iloc[1, 3], ipca["12m"])  # 12M
final_onshore["df_final_CarteiraBenchmark"].iloc[2, 4] = tax_diff(
    final_onshore["df_final_CarteiraBenchmark"].iloc[1, 4], ipca["24m"])  # 24M
final_onshore["df_final_CarteiraBenchmark"].iloc[2, 5] = tax_diff(
    final_onshore["df_final_CarteiraBenchmark"].iloc[1, 5], ipca["desde_ini"])  # DESDE O INICIO

final_onshore["df_final_CarteiraBenchmark"].iloc[0] = final_onshore["df_final_CarteiraBenchmark"].iloc[0].apply(
    lambda x: round_if_numeric(x, 0, "R$ "))
final_onshore["df_final_CarteiraBenchmark"].iloc[1:] = final_onshore["df_final_CarteiraBenchmark"].iloc[1:].applymap(
    lambda x: round_if_numeric(x, 1, "%", False))

final_offshore["df_final_CarteiraBenchmark"] = cmd_offshore["df_CarteiraBenchmark"].iloc[:2, :].copy()
final_offshore["df_final_CarteiraBenchmark"].replace(
    'nd', '-', inplace=True)
final_offshore["df_final_CarteiraBenchmark"].iloc[0] = final_offshore["df_final_CarteiraBenchmark"].iloc[0].apply(
    lambda x: round_if_numeric(x, 0, "$ "))

final_offshore["df_final_CarteiraBenchmark"].loc[len(final_offshore["df_final_CarteiraBenchmark"].index)] = [
    "EFFR",
    accum_return(df_benchmarks["US_EFFR"], previous_month_last_workday(final_offshore["end_date"]).strftime(
        "%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1],  # Mês
    accum_return(df_benchmarks["US_EFFR"], previous_month_last_workday(datetime(final_offshore["end_date"].year, 1, 10)).strftime(
        "%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1],  # Ano
    accum_return(df_benchmarks["US_EFFR"], (final_offshore["end_date"]+relativedelta(
        months=-12)).strftime("%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1],  # 12m
    accum_return(df_benchmarks["US_EFFR"], (final_offshore["end_date"]+relativedelta(
        months=-24)).strftime("%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1],  # 24m
    accum_return(df_benchmarks["US_EFFR"], final_offshore["start_date"].strftime("%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1]]  # Desde inicio

final_offshore["df_final_CarteiraBenchmark"].loc[len(final_offshore["df_final_CarteiraBenchmark"].index)] = [
    "BBG BONDS",
    accum_return(df_benchmarks["GLAG"], previous_month_last_workday(final_offshore["end_date"]).strftime(
        "%Y-%m-%d"), final_onshore["end_date"].strftime("%Y-%m-%d"))[-1],  # Mês
    accum_return(df_benchmarks["GLAG"], previous_month_last_workday(datetime(final_offshore["end_date"].year, 1, 10)).strftime(
        "%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1],  # Ano
    accum_return(df_benchmarks["GLAG"], (final_offshore["end_date"]+relativedelta(
        months=-12)).strftime("%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1],  # 12m
    accum_return(df_benchmarks["GLAG"], (final_offshore["end_date"]+relativedelta(
        months=-24)).strftime("%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1],  # 24m
    accum_return(df_benchmarks["GLAG"], final_offshore["start_date"].strftime("%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1]]  # Desde inicio

final_offshore["df_final_CarteiraBenchmark"].loc[len(final_offshore["df_final_CarteiraBenchmark"].index)] = [
    "S&P",
    accum_return(df_benchmarks["SPX"], previous_month_last_workday(final_offshore["end_date"]).strftime(
        "%Y-%m-%d"), final_onshore["end_date"].strftime("%Y-%m-%d"))[-1],  # Mês
    accum_return(df_benchmarks["SPX"], previous_month_last_workday(datetime(final_offshore["end_date"].year, 1, 10)).strftime(
        "%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1],  # Ano
    accum_return(df_benchmarks["SPX"], (final_offshore["end_date"]+relativedelta(months=-12)).strftime(
        "%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1],  # 12m
    accum_return(df_benchmarks["SPX"], (final_offshore["end_date"]+relativedelta(months=-24)).strftime(
        "%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1],  # 24m
    accum_return(df_benchmarks["SPX"], final_offshore["start_date"].strftime("%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1]]  # Desde inicio

final_offshore["df_final_CarteiraBenchmark"].iloc[3,
                                                  :] = final_offshore["df_final_CarteiraBenchmark"].iloc[2, :]
final_offshore["df_final_CarteiraBenchmark"].iloc[2,
                                                  0] = "Retorno Real"
final_offshore["df_final_CarteiraBenchmark"].iloc[2, 1] = tax_diff(
    final_offshore["df_final_CarteiraBenchmark"].iloc[1, 1], cpi["mes"])  # MÊS
final_offshore["df_final_CarteiraBenchmark"].iloc[2, 2] = tax_diff(
    final_offshore["df_final_CarteiraBenchmark"].iloc[1, 2], cpi["ano"])  # ANO
final_offshore["df_final_CarteiraBenchmark"].iloc[2, 3] = tax_diff(
    final_offshore["df_final_CarteiraBenchmark"].iloc[1, 3], cpi["12m"])  # 12M
final_offshore["df_final_CarteiraBenchmark"].iloc[2, 4] = tax_diff(
    final_offshore["df_final_CarteiraBenchmark"].iloc[1, 4], cpi["24m"])  # 24M
final_offshore["df_final_CarteiraBenchmark"].iloc[2, 5] = tax_diff(
    final_offshore["df_final_CarteiraBenchmark"].iloc[1, 5], cpi["desde_ini"])  # DESDE O INICIO

final_offshore["df_final_CarteiraBenchmark"].iloc[1:] = final_offshore["df_final_CarteiraBenchmark"].iloc[1:].applymap(
    lambda x: round_if_numeric(x, 1, "%", False))

cmd_offshore["df_PL"]["USDBRL"] = yahooFinance.Ticker("BRL=X").history(start=cmd_offshore["df_PL"].index[0].strftime(
    "%Y-%m-%d"), end=(cmd_offshore["df_PL"].index[-1]+timedelta(1)).strftime("%Y-%m-%d"))["Close"].tz_localize(None)  # USDBRL
cmd_offshore["df_PL"]["PL BRL"] = cmd_offshore["df_PL"]["USDBRL"] * \
    cmd_offshore["df_PL"]["PL"]

final_onshore["PL_final"] = cmd_onshore["df_PL"]["PL"][-1]
final_offshore["PL_final"] = cmd_offshore["df_PL"]["PL BRL"][-1]

union_index = cmd_onshore["df_PL"]["PL"].index.union(
    cmd_offshore["df_PL"]["PL BRL"].index)

s1_reindexed = cmd_onshore["df_PL"]["PL"].ffill().reindex(
    union_index, fill_value=0)
s2_reindexed = cmd_offshore["df_PL"]["PL BRL"].ffill().reindex(
    union_index, fill_value=0)

PL_total = s1_reindexed + s2_reindexed

# SLIDE 2
final_onshore["df_final_FluxoResgate"] = pd.DataFrame({"Tempo de Resgate": ["D0 a D5", "D6 a D30", "D31+", "N/A"],
                                                        "Valor líquido": [cmd_onshore["df_FluxoResgate"].iloc[:2, :].sum()[1], cmd_onshore["df_FluxoResgate"].iloc[2, 1], cmd_onshore["df_FluxoResgate"].iloc[3, 1], cmd_onshore["df_FluxoResgate"].iloc[4, 1]],
                                                        "%": [cmd_onshore["df_FluxoResgate"].iloc[:2, :].sum()[2], cmd_onshore["df_FluxoResgate"].iloc[2, 2], cmd_onshore["df_FluxoResgate"].iloc[3, 2], cmd_onshore["df_FluxoResgate"].iloc[4, 2]]})
final_onshore["df_final_FluxoResgate"]["Valor líquido"] = final_onshore["df_final_FluxoResgate"]["Valor líquido"].apply(
    lambda x: round_if_numeric(x, 0, "R$ "))

final_onshore["df_final_RendimentoCarteira"] = pd.DataFrame({"Carteira": cmd_onshore["df_PL"]["Cota"],
                                                              # "CDI":df_CDI["rendimento"],
                                                              "PL": cmd_onshore["df_PL"]["PL"]})
final_onshore["df_final_RendimentoCarteira"]["CDI"] = accum_return(
    df_benchmarks.loc[:, "CDI"], final_onshore["df_final_RendimentoCarteira"].index[0], final_onshore["df_final_RendimentoCarteira"].index[-1]).ffill()
final_onshore["df_final_RendimentoCarteira"] = final_onshore["df_final_RendimentoCarteira"][[
    "Carteira", "CDI", "PL"]]
final_onshore["df_final_RendimentoCarteira"].reset_index(
    inplace=True)

# SLIDE 3
final_onshore["df_final_RetornosNominais"] = cmd_onshore["df_RetornosNominais"].copy(
)
final_onshore["df_final_RetornosNominais"].rename(
    columns={final_onshore["df_final_RetornosNominais"].columns[0]: "Ano_ref"}, inplace=True)
final_onshore["df_final_RetornosNominais"] = final_onshore["df_final_RetornosNominais"].sort_values(
    by='Ativo', ascending=False)
final_onshore["df_final_RetornosNominais"] = final_onshore["df_final_RetornosNominais"].sort_values(
    by='Ano_ref', ascending=False)
final_onshore["df_final_RetornosNominais"] = final_onshore["df_final_RetornosNominais"].reset_index(
    drop=True)
final_onshore["df_final_RetornosNominais"] = final_onshore["df_final_RetornosNominais"].replace(
    "", "-")
final_onshore["df_final_RetornosNominais"].iloc[:, 2:] = final_onshore["df_final_RetornosNominais"].iloc[:, 2:].applymap(
    lambda x: round_if_numeric(x, 1, "%", False))

final_onshore["df_final_ResumoPorConta"] = cmd_onshore["df_ResumoPorConta"].copy(
).reset_index(drop=True)
final_onshore["df_final_ResumoPorConta"]["Ativo"] = final_onshore["df_final_ResumoPorConta"]["Ativo"].str.upper(
)
final_onshore["df_final_ResumoPorConta"].iloc[:, 2:-1] = final_onshore["df_final_ResumoPorConta"].iloc[:,
                                                                                                        2:-1].applymap(lambda x: round_if_numeric(x, 0, "R$ "))
final_onshore["df_final_ResumoPorConta"].iloc[:, -1] = final_onshore["df_final_ResumoPorConta"].iloc[:, -1].apply(
    lambda x: round_if_numeric(x, 1, "%", False))

# SLIDE 4
try:
    final_onshore["df_final_PerfAttr"] = cmd_onshore["df_PerfAttr"][["Classe", "Bps no Mês"]].replace('', np.nan).dropna().sort_values(
        by=['Bps no Mês'], ascending=False)
except:
    final_onshore["df_final_PerfAttr"] = cmd_onshore["df_PerfAttr"][["Classe", "Bps no Período"]].replace('', np.nan).dropna().sort_values(
        by=['Bps no Período'], ascending=False)
    final_onshore["df_final_PerfAttr"].rename(
        columns={'Bps no Período': 'Bps no Mês'}, inplace=True)

final_onshore["df_final_PerfAttr"] = final_onshore["df_final_PerfAttr"][final_onshore["df_final_PerfAttr"]["Bps no Mês"] != 0].reset_index(
    drop=True)
final_onshore["df_final_PerfAttr"] = pd.concat([final_onshore["df_final_PerfAttr"].drop(final_onshore["df_final_PerfAttr"].index[0]), pd.DataFrame(
    [final_onshore["df_final_PerfAttr"].iloc[0]])]).reset_index(drop=True)

final_onshore["df_final_AlocacaoClasse"] = cmd_onshore["df_ClasseDeAtivos"][[
    "Classe", "% da Carteira"]].sort_values(by=['% da Carteira'], ascending=False)

final_onshore["df_final_Volatilidade"] = cmd_onshore["df_Volatilidade"][cmd_onshore["df_Volatilidade"].Ativo.isin(
    final_onshore["df_final_PerfAttr"].Classe)]
final_onshore["df_final_Volatilidade"] = final_onshore["df_final_Volatilidade"][[
    "Ativo", "Vol. 12m"]]

# SLIDE 6

final_onshore["df_final_RetornoNominalAtivo"] = cmd_onshore["df_RetornoNominalAtivo"].reset_index(
    drop=True)
final_onshore["df_final_RetornoNominalAtivo"] = final_onshore["df_final_RetornoNominalAtivo"][final_onshore["df_final_RetornoNominalAtivo"]["Ativo"] != "% do CDI"].reset_index(
    drop=True)
final_onshore["df_final_RetornoNominalAtivo"]["Data Compra"] = pd.to_datetime(
    final_onshore["df_final_RetornoNominalAtivo"]["Data Compra"], dayfirst=True).dt.date
final_onshore["df_final_RetornoNominalAtivo"]["R$ Alocado"] = final_onshore[
    "df_final_RetornoNominalAtivo"]["% Alocado"]*cmd_onshore["df_PL"]["PL"][-1]/100

class_index = final_onshore["df_final_RetornoNominalAtivo"].index[final_onshore["df_final_RetornoNominalAtivo"]['Ativo'].isin(
    final_onshore["df_final_AlocacaoClasse"].Classe)].tolist()
class_index.append(
    len(final_onshore["df_final_RetornoNominalAtivo"]))
final_onshore["dfs_class"] = {}
for i in range(len(class_index)-1):
    categoria_atual = final_onshore["df_final_RetornoNominalAtivo"].iloc[class_index[i]]['Ativo']
    final_onshore["dfs_class"][categoria_atual] = final_onshore[
        "df_final_RetornoNominalAtivo"].iloc[class_index[i]:class_index[i+1]]

for classe in list(final_onshore["dfs_class"]):
    try:
        final_onshore["dfs_class"][classe].reset_index(drop=True)
        benchmark_code = classe_benchmark[classe_benchmark["class_region"]
                                          == (classe+"_Onshore")]["benchmark_code"].iloc[0]
        benchmark_name = classe_benchmark[classe_benchmark["class_region"]
                                          == (classe+"_Onshore")]["benchmark_name"].iloc[0]
        benchmark_metric = classe_benchmark[classe_benchmark["class_region"]
                                            == (classe+"_Onshore")]["metric"].iloc[0]
        benchmark_name = f"% do {benchmark_name}" if benchmark_metric == "%" else f"Excesso sobre {benchmark_name}"
        benchmark = df_benchmarks[benchmark_code]

        counter = 1
        for i in range(1, len(final_onshore["dfs_class"][classe])):
            ativo = final_onshore["dfs_class"][classe].iloc[i+counter-1]
            line = pd.DataFrame({'Ativo': benchmark_name,
                                  'Mês': benchmark_diff(ativo["Mês"],
                                                        accum_return(benchmark, previous_month_last_workday(final_onshore["end_date"]).strftime(
                                                            "%Y-%m-%d"), final_onshore["end_date"].strftime("%Y-%m-%d"))[-1], benchmark_metric),
                                  'Ano': benchmark_diff(ativo["Ano"], accum_return(benchmark, previous_month_last_workday(datetime(final_offshore["end_date"].year, 1, 10)).strftime(
                                      "%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1], benchmark_metric),
                                  '03meses': benchmark_diff(ativo["03meses"],
                                                            accum_return(benchmark, (final_offshore["end_date"]+relativedelta(months=-3)).strftime(
                                                                "%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1], benchmark_metric),
                                  '06meses': benchmark_diff(ativo["06meses"],
                                                            accum_return(benchmark, (final_offshore["end_date"]+relativedelta(months=-6)).strftime(
                                                                "%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1], benchmark_metric),
                                  '12meses': benchmark_diff(ativo["12meses"],
                                                            accum_return(benchmark, (final_offshore["end_date"]+relativedelta(months=-12)).strftime(
                                                                "%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1], benchmark_metric),
                                  '24meses': benchmark_diff(ativo["24meses"],
                                                            accum_return(benchmark, (final_offshore["end_date"]+relativedelta(months=-24)).strftime(
                                                                "%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1], benchmark_metric),
                                  '36meses': benchmark_diff(ativo["36meses"],
                                                            accum_return(benchmark, (final_offshore["end_date"]+relativedelta(months=-36)).strftime(
                                                                "%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1], benchmark_metric),
                                  'Desde aCompra': benchmark_diff(ativo["Desde aCompra"],
                                                                  accum_return(benchmark, ativo["Data Compra"].strftime(
                                                                      "%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1], benchmark_metric),
                                  'Data Compra': "",
                                  '% Alocado': "",
                                  'R$ Alocado': ""}, index=[1000+i])
            final_onshore["dfs_class"][classe] = pd.concat([final_onshore["dfs_class"][classe].iloc[:i+counter], line, final_onshore["dfs_class"][classe].iloc[i+counter:]]
                                                            )
            counter += 1
        final_onshore["dfs_class"][classe] = final_onshore["dfs_class"][classe][["Ativo", "% Alocado", "R$ Alocado", "Mês", "Ano", '03meses', '06meses', '12meses', '24meses',
                                                                                  '36meses', 'Desde aCompra']]

        final_onshore["dfs_class"][classe] = final_onshore["dfs_class"][classe].replace(
            "", "-")
        final_onshore["dfs_class"][classe] = final_onshore["dfs_class"][classe].replace(
            "nd", "-")

        final_onshore["dfs_class"][classe].iloc[:, 1] = final_onshore["dfs_class"][classe].iloc[:, 1].apply(
            lambda x: round_if_numeric(x, 1, "%", False))

        final_onshore["dfs_class"][classe].iloc[:, 3:] = final_onshore["dfs_class"][classe].iloc[:, 3:].applymap(
            lambda x: round_if_numeric(x, 1, "%", False))

        final_onshore["dfs_class"][classe].iloc[:, 2] = final_onshore["dfs_class"][classe].iloc[:, 2].apply(
            lambda x: round_if_numeric(x, 0, "R$", True))

    except:
        print(f"     Error on ONSHORE class: {classe}")
        del final_onshore["dfs_class"][classe]

############ OFFSHORE ###########################
# SLIDE 7
final_offshore["df_final_FluxoResgate"] = pd.DataFrame({"Tempo de Resgate": ["D0 a D5", "D6 a D30", "D31+", "N/A"],
                                                        "Valor líquido": [cmd_offshore["df_FluxoResgate"].iloc[:2, :].sum()[1], cmd_offshore["df_FluxoResgate"].iloc[2, 1], cmd_offshore["df_FluxoResgate"].iloc[3, 1], cmd_offshore["df_FluxoResgate"].iloc[4, 1]],
                                                        "%": [cmd_offshore["df_FluxoResgate"].iloc[:2, :].sum()[2], cmd_offshore["df_FluxoResgate"].iloc[2, 2], cmd_offshore["df_FluxoResgate"].iloc[3, 2], cmd_offshore["df_FluxoResgate"].iloc[4, 2]]})
final_offshore["df_final_FluxoResgate"]["Valor líquido"] = final_offshore["df_final_FluxoResgate"]["Valor líquido"].apply(
    lambda x: round_if_numeric(x, 0, "$ "))

df_final_RendimentoCarteira = pd.DataFrame({"Carteira": cmd_offshore["df_PL"]["Cota"],
                                            # "CDI":df_CDI["rendimento"],
                                            "PL": cmd_offshore["df_PL"]["PL"]})

final_offshore["df_final_RendimentoCarteira"] = pd.DataFrame({"Carteira": cmd_offshore["df_PL"]["Cota"],
                                                              # "CDI":df_CDI["rendimento"],
                                                              "PL": cmd_offshore["df_PL"]["PL"]})
final_offshore["df_final_RendimentoCarteira"]["EFFR"] = accum_return(
    df_benchmarks.loc[:, "US_EFFR"], final_offshore["df_final_RendimentoCarteira"].index[0], final_offshore["df_final_RendimentoCarteira"].index[-1]).ffill()
final_offshore["df_final_RendimentoCarteira"] = final_offshore["df_final_RendimentoCarteira"][[
    "Carteira", "EFFR", "PL"]]
final_offshore["df_final_RendimentoCarteira"].reset_index(
    inplace=True)

# SLIDE 8
final_offshore["df_final_RetornosNominais"] = cmd_offshore["df_RetornosNominais"].copy(
)
final_offshore["df_final_RetornosNominais"].iloc[:, 1] = final_offshore["df_final_RetornosNominais"].iloc[:, 1].apply(
    lambda x: "%Fed Funds" if "EFFR" in x else x)
final_offshore["df_final_RetornosNominais"].rename(
    columns={final_offshore["df_final_RetornosNominais"].columns[0]: "Ano_ref"}, inplace=True)
final_offshore["df_final_RetornosNominais"] = final_offshore["df_final_RetornosNominais"].sort_values(
    by='Ano_ref', ascending=False)
final_offshore["df_final_RetornosNominais"] = final_offshore["df_final_RetornosNominais"].reset_index(
    drop=True)
final_offshore["df_final_RetornosNominais"] = final_offshore["df_final_RetornosNominais"].replace(
    "", "-")
final_offshore["df_final_RetornosNominais"].iloc[:, 2:] = final_offshore["df_final_RetornosNominais"].iloc[:, 2:].applymap(
    lambda x: round_if_numeric(x, 1, "%", False))

final_offshore["df_final_ResumoPorConta"] = cmd_offshore["df_ResumoPorConta"].copy(
).reset_index(drop=True)
final_offshore["df_final_ResumoPorConta"]["Ativo"] = final_offshore["df_final_ResumoPorConta"]["Ativo"].str.upper(
)
final_offshore["df_final_ResumoPorConta"].iloc[:, 2:-1] = final_offshore["df_final_ResumoPorConta"].iloc[:,
                                                                                                          2:-1].applymap(lambda x: round_if_numeric(x, 0, "$ "))
final_offshore["df_final_ResumoPorConta"].iloc[:, -1] = final_offshore["df_final_ResumoPorConta"].iloc[:, -1].apply(
    lambda x: round_if_numeric(x, 1, "%", False))

# SLIDE 9
try:
    final_offshore["df_final_PerfAttr"] = cmd_offshore["df_PerfAttr"][["Classe", "Bps no Mês"]].replace("", np.nan).sort_values(
        by=['Bps no Mês'], ascending=False)
except:
    final_offshore["df_final_PerfAttr"] = cmd_offshore["df_PerfAttr"][["Classe", "Bps no Período"]].replace("", np.nan).sort_values(
        by=['Bps no Período'], ascending=False)
    final_offshore["df_final_PerfAttr"].rename(
        columns={'Bps no Período': 'Bps no Mês'}, inplace=True)

final_offshore["df_final_PerfAttr"] = final_offshore["df_final_PerfAttr"][final_offshore["df_final_PerfAttr"]["Bps no Mês"] != 0].reset_index(
    drop=True)
final_offshore["df_final_PerfAttr"] = pd.concat([final_offshore["df_final_PerfAttr"].drop(final_offshore["df_final_PerfAttr"].index[0]), pd.DataFrame(
    [final_offshore["df_final_PerfAttr"].iloc[0]])]).reset_index(drop=True)

final_offshore["df_final_AlocacaoClasse"] = cmd_offshore["df_ClasseDeAtivos"][[
    "Classe", "% da Carteira"]].sort_values(by=['% da Carteira'], ascending=False)

final_offshore["df_final_Volatilidade"] = cmd_offshore["df_Volatilidade"][cmd_offshore["df_Volatilidade"].Ativo.isin(
    final_offshore["df_final_PerfAttr"].Classe)]
final_offshore["df_final_Volatilidade"] = final_offshore["df_final_Volatilidade"][[
    "Ativo", "Vol. 12m"]]

# SLIDE TABELA BONI OFFSHORE
carteira_ekho_offshore_mes = pd.merge(carteira_ekho_offshore, cmd_offshore["df_PerfAttr"].iloc[:,1:7], left_on="classe", right_on="Classe", how="left")
carteira_ekho_offshore_mes = carteira_ekho_offshore_mes.drop(columns="Classe")
carteira_ekho_offshore_mes = carteira_ekho_offshore_mes.fillna(0)
carteira_ekho_offshore_mes["benchmark_return"] =  carteira_ekho_offshore_mes["benchmark"].apply(lambda x: accum_return(df_benchmarks[x],previous_month_last_workday(final_offshore["end_date"]).strftime(
    "%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1])
carteira_ekho_offshore_mes = carteira_ekho_offshore_mes.replace("nd",0)

column_indices = [0,1,2,3,4,5,6,7,8,9]
new_names = ['classe','peso_pl','benchmark','benchmark_name','saldo_bruto','peso_aloc','performance','contrib_bruta_value','contrib_bruta_percent','perf_benchmark']
old_names = carteira_ekho_offshore_mes.columns[column_indices]
carteira_ekho_offshore_mes.rename(columns=dict(zip(old_names, new_names)), inplace=True)
carteira_ekho_offshore_mes['contrib_bruta_percent'] = carteira_ekho_offshore_mes['contrib_bruta_percent']/100
carteira_ekho_offshore_mes["peso_pl"] = carteira_ekho_offshore_mes["peso_pl"]*100

carteira_ekho_offshore_mes["perf_vs_bench"] = carteira_ekho_offshore_mes["performance"]/carteira_ekho_offshore_mes["perf_benchmark"]*100
carteira_ekho_offshore_mes["contrib_bench"] = (carteira_ekho_offshore_mes["peso_pl"]/100*carteira_ekho_offshore_mes["perf_benchmark"]/100)*100
carteira_ekho_offshore_mes["total"] = carteira_ekho_offshore_mes["contrib_bruta_percent"]-carteira_ekho_offshore_mes["contrib_bench"]
carteira_ekho_offshore_mes["alocacao"] = (((carteira_ekho_offshore_mes["peso_aloc"]/100)-(carteira_ekho_offshore_mes["peso_pl"]/100))*carteira_ekho_offshore_mes["perf_benchmark"]/100)*100
carteira_ekho_offshore_mes["selecao"] = carteira_ekho_offshore_mes["total"] - carteira_ekho_offshore_mes["alocacao"]

carteira_ekho_offshore_mes = carteira_ekho_offshore_mes[["classe","peso_aloc", "saldo_bruto", "performance", "perf_vs_bench", "contrib_bruta_percent", "contrib_bruta_value", "peso_pl","perf_benchmark", "benchmark_name", "contrib_bench", "total", "alocacao", "selecao"]]
carteira_ekho_offshore_mes = carteira_ekho_offshore_mes.fillna(0)


carteira_ekho_offshore_ano = pd.merge(carteira_ekho_offshore, cmd_offshore["df_PerfAttr"].iloc[:,[1,2,3,7,8,9]], left_on="classe", right_on="Classe", how="left")
carteira_ekho_offshore_ano = carteira_ekho_offshore_ano.drop(columns="Classe")
carteira_ekho_offshore_ano = carteira_ekho_offshore_ano.fillna(0)
carteira_ekho_offshore_ano["benchmark_return"] =  carteira_ekho_offshore_ano["benchmark"].apply(lambda x: accum_return(df_benchmarks[x],previous_month_last_workday(datetime(final_offshore["end_date"].year, 1, 10)).strftime(
    "%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1])
carteira_ekho_offshore_ano = carteira_ekho_offshore_ano.replace("nd",0)

column_indices = [0,1,2,3,4,5,6,7,8,9]
new_names = ['classe','peso_pl','benchmark','benchmark_name','saldo_bruto','peso_aloc','performance','contrib_bruta_value','contrib_bruta_percent','perf_benchmark']
old_names = carteira_ekho_offshore_ano.columns[column_indices]
carteira_ekho_offshore_ano.rename(columns=dict(zip(old_names, new_names)), inplace=True)
carteira_ekho_offshore_ano['contrib_bruta_percent'] = carteira_ekho_offshore_ano['contrib_bruta_percent']/100
carteira_ekho_offshore_ano["peso_pl"] = carteira_ekho_offshore_ano["peso_pl"]*100

carteira_ekho_offshore_ano["perf_vs_bench"] = carteira_ekho_offshore_ano["performance"]/carteira_ekho_offshore_ano["perf_benchmark"]*100
carteira_ekho_offshore_ano["contrib_bench"] = (carteira_ekho_offshore_ano["peso_pl"]/100*carteira_ekho_offshore_ano["perf_benchmark"]/100)*100
carteira_ekho_offshore_ano["total"] = carteira_ekho_offshore_ano["contrib_bruta_percent"]-carteira_ekho_offshore_ano["contrib_bench"]
carteira_ekho_offshore_ano["alocacao"] = (((carteira_ekho_offshore_ano["peso_aloc"]/100)-(carteira_ekho_offshore_ano["peso_pl"]/100))*carteira_ekho_offshore_ano["perf_benchmark"]/100)*100
carteira_ekho_offshore_ano["selecao"] = carteira_ekho_offshore_ano["total"] - carteira_ekho_offshore_ano["alocacao"]

carteira_ekho_offshore_ano = carteira_ekho_offshore_ano[["classe","peso_aloc", "saldo_bruto", "performance", "perf_vs_bench", "contrib_bruta_percent", "contrib_bruta_value", "peso_pl","perf_benchmark", "benchmark_name", "contrib_bench", "total", "alocacao", "selecao"]]
carteira_ekho_offshore_ano = carteira_ekho_offshore_ano.fillna(0)

# SLIDE TABELA BONI ONSHORE
carteira_ekho_onshore_mes = pd.merge(carteira_ekho_onshore, cmd_onshore["df_PerfAttr"].iloc[:,1:7], left_on="classe", right_on="Classe", how="left")
carteira_ekho_onshore_mes = carteira_ekho_onshore_mes.drop(columns="Classe")
carteira_ekho_onshore_mes = carteira_ekho_onshore_mes.fillna(0)
carteira_ekho_onshore_mes["benchmark_return"] =  carteira_ekho_onshore_mes["benchmark"].apply(lambda x: accum_return(df_benchmarks[x],previous_month_last_workday(final_onshore["end_date"]).strftime(
    "%Y-%m-%d"), final_onshore["end_date"].strftime("%Y-%m-%d"))[-1])
carteira_ekho_onshore_mes = carteira_ekho_onshore_mes.replace("nd",0)

column_indices = [0,1,2,3,4,5,6,7,8,9]
new_names = ['classe','peso_pl','benchmark','benchmark_name','saldo_bruto','peso_aloc','performance','contrib_bruta_value','contrib_bruta_percent','perf_benchmark']
old_names = carteira_ekho_onshore_mes.columns[column_indices]
carteira_ekho_onshore_mes.rename(columns=dict(zip(old_names, new_names)), inplace=True)
carteira_ekho_onshore_mes['contrib_bruta_percent'] = carteira_ekho_onshore_mes['contrib_bruta_percent']/100
carteira_ekho_onshore_mes["peso_pl"] = carteira_ekho_onshore_mes["peso_pl"]*100

carteira_ekho_onshore_mes["perf_vs_bench"] = carteira_ekho_onshore_mes["performance"]/carteira_ekho_onshore_mes["perf_benchmark"]*100
carteira_ekho_onshore_mes["contrib_bench"] = (carteira_ekho_onshore_mes["peso_pl"]/100*carteira_ekho_onshore_mes["perf_benchmark"]/100)*100
carteira_ekho_onshore_mes["total"] = carteira_ekho_onshore_mes["contrib_bruta_percent"]-carteira_ekho_onshore_mes["contrib_bench"]
carteira_ekho_onshore_mes["alocacao"] = (((carteira_ekho_onshore_mes["peso_aloc"]/100)-(carteira_ekho_onshore_mes["peso_pl"]/100))*carteira_ekho_onshore_mes["perf_benchmark"]/100)*100
carteira_ekho_onshore_mes["selecao"] = carteira_ekho_onshore_mes["total"] - carteira_ekho_onshore_mes["alocacao"]

carteira_ekho_onshore_mes = carteira_ekho_onshore_mes[["classe","peso_aloc", "saldo_bruto", "performance", "perf_vs_bench", "contrib_bruta_percent", "contrib_bruta_value", "peso_pl","perf_benchmark", "benchmark_name", "contrib_bench", "total", "alocacao", "selecao"]]
carteira_ekho_onshore_mes = carteira_ekho_onshore_mes.fillna(0)


carteira_ekho_onshore_ano = pd.merge(carteira_ekho_onshore, cmd_onshore["df_PerfAttr"].iloc[:,[1,2,3,7,8,9]], left_on="classe", right_on="Classe", how="left")
carteira_ekho_onshore_ano = carteira_ekho_onshore_ano.drop(columns="Classe")
carteira_ekho_onshore_ano = carteira_ekho_onshore_ano.fillna(0)
carteira_ekho_onshore_ano["benchmark_return"] =  carteira_ekho_onshore_ano["benchmark"].apply(lambda x: accum_return(df_benchmarks[x],previous_month_last_workday(datetime(final_offshore["end_date"].year, 1, 10)).strftime(
    "%Y-%m-%d"), final_onshore["end_date"].strftime("%Y-%m-%d"))[-1])
carteira_ekho_onshore_ano = carteira_ekho_onshore_ano.replace("nd",0)

column_indices = [0,1,2,3,4,5,6,7,8,9]
new_names = ['classe','peso_pl','benchmark','benchmark_name','saldo_bruto','peso_aloc','performance','contrib_bruta_value','contrib_bruta_percent','perf_benchmark']
old_names = carteira_ekho_onshore_ano.columns[column_indices]
carteira_ekho_onshore_ano.rename(columns=dict(zip(old_names, new_names)), inplace=True)
carteira_ekho_onshore_ano['contrib_bruta_percent'] = carteira_ekho_onshore_ano['contrib_bruta_percent']/100
carteira_ekho_onshore_ano["peso_pl"] = carteira_ekho_onshore_ano["peso_pl"]*100

carteira_ekho_onshore_ano["perf_vs_bench"] = carteira_ekho_onshore_ano["performance"]/carteira_ekho_onshore_ano["perf_benchmark"]*100
carteira_ekho_onshore_ano["contrib_bench"] = (carteira_ekho_onshore_ano["peso_pl"]/100*carteira_ekho_onshore_ano["perf_benchmark"]/100)*100
carteira_ekho_onshore_ano["total"] = carteira_ekho_onshore_ano["contrib_bruta_percent"]-carteira_ekho_onshore_ano["contrib_bench"]
carteira_ekho_onshore_ano["alocacao"] = (((carteira_ekho_onshore_ano["peso_aloc"]/100)-(carteira_ekho_onshore_ano["peso_pl"]/100))*carteira_ekho_onshore_ano["perf_benchmark"]/100)*100
carteira_ekho_onshore_ano["selecao"] = carteira_ekho_onshore_ano["total"] - carteira_ekho_onshore_ano["alocacao"]

carteira_ekho_onshore_ano = carteira_ekho_onshore_ano[["classe","peso_aloc", "saldo_bruto", "performance", "perf_vs_bench", "contrib_bruta_percent", "contrib_bruta_value", "peso_pl","perf_benchmark", "benchmark_name", "contrib_bench", "total", "alocacao", "selecao"]]
carteira_ekho_onshore_ano = carteira_ekho_onshore_ano.fillna(0)

# SLIDE 10

final_offshore["df_final_RetornoNominalAtivo"] = cmd_offshore["df_RetornoNominalAtivo"].reset_index(
    drop=True)
final_offshore["df_final_RetornoNominalAtivo"] = final_offshore["df_final_RetornoNominalAtivo"][final_offshore["df_final_RetornoNominalAtivo"]["Ativo"] != "% do CDI"].reset_index(
    drop=True)
final_offshore["df_final_RetornoNominalAtivo"]["Data Compra"] = pd.to_datetime(
    final_offshore["df_final_RetornoNominalAtivo"]["Data Compra"], dayfirst=True).dt.date
final_offshore["df_final_RetornoNominalAtivo"]["$ Alocado"] = final_offshore[
    "df_final_RetornoNominalAtivo"]["% Alocado"]*cmd_offshore["df_PL"]["PL"][-1]/100

final_offshore["df_final_RetornoNominalAtivo"] = final_offshore["df_final_RetornoNominalAtivo"].replace(
    0, "nd")
class_index = final_offshore["df_final_RetornoNominalAtivo"].index[final_offshore["df_final_RetornoNominalAtivo"]['Ativo'].isin(
    final_offshore["df_final_AlocacaoClasse"].Classe)].tolist()
class_index.append(
    len(final_offshore["df_final_RetornoNominalAtivo"]))
final_offshore["dfs_class"] = {}
for i in range(len(class_index)-1):
    categoria_atual = final_offshore["df_final_RetornoNominalAtivo"].iloc[class_index[i]]['Ativo']
    final_offshore["dfs_class"][categoria_atual] = final_offshore[
        "df_final_RetornoNominalAtivo"].iloc[class_index[i]:class_index[i+1]]

for classe in list(final_offshore["dfs_class"]):
    try:
        final_offshore["dfs_class"][classe].reset_index(drop=True)
        benchmark_code = classe_benchmark[classe_benchmark["class_region"]
                                          == (classe+"_Offshore")]["benchmark_code"].iloc[0]
        benchmark_name = classe_benchmark[classe_benchmark["class_region"]
                                          == (classe+"_Offshore")]["benchmark_name"].iloc[0]
        benchmark_metric = classe_benchmark[classe_benchmark["class_region"]
                                            == (classe+"_Offshore")]["metric"].iloc[0]
        benchmark_name = f"% do {benchmark_name}" if benchmark_metric == "%" else f"Excesso sobre {benchmark_name}"
        benchmark = df_benchmarks[benchmark_code]

        counter = 1
        for i in range(1, len(final_offshore["dfs_class"][classe])):
            ativo = final_offshore["dfs_class"][classe].iloc[i+counter-1]
            line = pd.DataFrame({'Ativo': benchmark_name,
                                  'Mês': benchmark_diff(ativo["Mês"],
                                                        accum_return(benchmark, previous_month_last_workday(final_offshore["end_date"]).strftime(
                                                            "%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1], benchmark_metric),
                                  'Ano': benchmark_diff(ativo["Ano"], accum_return(benchmark, previous_month_last_workday(datetime(final_offshore["end_date"].year, 1, 10)).strftime(
                                      "%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1], benchmark_metric),
                                  '03meses': benchmark_diff(ativo["03meses"],
                                                            accum_return(benchmark, (final_offshore["end_date"]+relativedelta(months=-3)).strftime(
                                                                "%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1], benchmark_metric),
                                  '06meses': benchmark_diff(ativo["06meses"],
                                                            accum_return(benchmark, (final_offshore["end_date"]+relativedelta(months=-6)).strftime(
                                                                "%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1], benchmark_metric),
                                  '12meses': benchmark_diff(ativo["12meses"],
                                                            accum_return(benchmark, (final_offshore["end_date"]+relativedelta(months=-12)).strftime(
                                                                "%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1], benchmark_metric),
                                  '24meses': benchmark_diff(ativo["24meses"],
                                                            accum_return(benchmark, (final_offshore["end_date"]+relativedelta(months=-24)).strftime(
                                                                "%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1], benchmark_metric),
                                  '36meses': benchmark_diff(ativo["36meses"],
                                                            accum_return(benchmark, (final_offshore["end_date"]+relativedelta(months=-36)).strftime(
                                                                "%Y-%m-%d"), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1], benchmark_metric),
                                  'Desde aCompra': benchmark_diff(ativo["Desde aCompra"],
                                                                  accum_return(benchmark, ("2023-01-02" if pd.isnull(ativo["Data Compra"]) else ativo["Data Compra"].strftime(
                                                                      "%Y-%m-%d")), final_offshore["end_date"].strftime("%Y-%m-%d"))[-1], benchmark_metric),
                                  'Data Compra': "",
                                  '% Alocado': "",
                                  '$ Alocado': ""}, index=[1000+i])
            final_offshore["dfs_class"][classe] = pd.concat([final_offshore["dfs_class"][classe].iloc[:i+counter], line, final_offshore["dfs_class"][classe].iloc[i+counter:]]
                                                            )
            counter += 1
        final_offshore["dfs_class"][classe] = final_offshore["dfs_class"][classe][["Ativo", "% Alocado", "$ Alocado", "Mês", "Ano", '03meses', '06meses', '12meses', '24meses',
                                                                                    '36meses', 'Desde aCompra']]

        final_offshore["dfs_class"][classe] = final_offshore["dfs_class"][classe].replace(
            "", "-")
        final_offshore["dfs_class"][classe] = final_offshore["dfs_class"][classe].replace(
            "nd", "-")

        final_offshore["dfs_class"][classe].iloc[:, 1] = final_offshore["dfs_class"][classe].iloc[:, 1].apply(
            lambda x: round_if_numeric(x, 1, "%", False))

        final_offshore["dfs_class"][classe].iloc[:, 3:] = final_offshore["dfs_class"][classe].iloc[:, 3:].applymap(
            lambda x: round_if_numeric(x, 1, "%", False))

        final_offshore["dfs_class"][classe].iloc[:, 2] = final_offshore["dfs_class"][classe].iloc[:, 2].apply(
            lambda x: round_if_numeric(x, 0, "$", True))
    except:
        print(f"     Error on OFFSHORE class: {classe}")
        del final_offshore["dfs_class"][classe]

  ######################## PYTHON PPTX #########################

prs = Presentation("templates/Template Relatório - ON + OFF.pptx")

# date
for slide_idx in range(len(prs.slides)):
    slide = prs.slides[slide_idx]
    try:
        date_shape = [
            shape for shape in slide.shapes if shape.name == "date"][0]
        date_shape.text_frame.paragraphs[0].runs[0].text = date_title
    except:
        pass

# name
for slide_idx in range(len(prs.slides)):
    slide = prs.slides[slide_idx]
    try:
        name_shape = [
            shape for shape in slide.shapes if shape.name == "name"][0]
        name_shape.text_frame.paragraphs[0].runs[0].text = client_name
    except:
        pass

# slide 1 - Resumo Geral
slide = prs.slides[1]
carteiraBenchmark_on_table = [
    shape for shape in slide.shapes if shape.name == "CarteiraBenchmark_on"][0].table
carteiraBenchmark_off_table = [
    shape for shape in slide.shapes if shape.name == "CarteiraBenchmark_off"][0].table
geografica_graf = [
    shape for shape in slide.shapes if shape.name == "geografica"][0].chart
evolucao_graf = [
    shape for shape in slide.shapes if shape.name == "evolucao"][0].chart
diversificacao_patrimonio_table = [
    shape for shape in slide.shapes if shape.name == "diversificacao_patrimonio"][0].table

# Onshore / Offshore
for i in range(5):
    for j in range(5):
        carteiraBenchmark_on_table.cell(
            i+2, j+1).text_frame.paragraphs[0].runs[0].text = final_onshore["df_final_CarteiraBenchmark"].iloc[i, j+1]
        carteiraBenchmark_off_table.cell(
            i+2, j+1).text_frame.paragraphs[0].runs[0].text = final_offshore["df_final_CarteiraBenchmark"].iloc[i, j+1]

# ---define new chart data---
chart_data = CategoryChartData()
chart_data.categories = [date for date in PL_total.index]
chart_data.add_series('PL', tuple(
    PL_total.ffill().fillna(0).values))
# ---replace chart data---
evolucao_graf.replace_data(chart_data)

# ---define new chart data---
chart_data = CategoryChartData()
chart_data.categories = ["ONSHORE", "OFFSHORE"]
chart_data.add_series(
    'PL', (final_onshore["PL_final"], final_offshore["PL_final"]))
# ---replace chart data---
geografica_graf.replace_data(chart_data)

diversificacao_patrimonio_table.cell(1, 1).text_frame.paragraphs[0].runs[0].text = round_if_numeric(
    final_onshore["PL_final"], 0, "R$ ")
diversificacao_patrimonio_table.cell(2, 1).text_frame.paragraphs[0].runs[0].text = round_if_numeric(
    final_offshore["PL_final"], 0, "R$ ")
diversificacao_patrimonio_table.cell(3, 1).text_frame.paragraphs[0].runs[0].text = round_if_numeric(
    final_onshore["PL_final"]+final_offshore["PL_final"], 0, "R$ ")

# slide 3 - Resumo Onshore
slide = prs.slides[3]
carteiraBenchmark_on_table = [
    shape for shape in slide.shapes if shape.name == "CarteiraBenchmark"][0].table
liq_val1 = [shape for shape in slide.shapes if shape.name ==
            "liq_val1"][0].text_frame.paragraphs[0].runs[0]
liq_val2 = [shape for shape in slide.shapes if shape.name ==
            "liq_val2"][0].text_frame.paragraphs[0].runs[0]
liq_val3 = [shape for shape in slide.shapes if shape.name ==
            "liq_val3"][0].text_frame.paragraphs[0].runs[0]
liq_val4 = [shape for shape in slide.shapes if shape.name ==
            "liq_val4"][0].text_frame.paragraphs[0].runs[0]
liq1 = [shape for shape in slide.shapes if shape.name == "liq1"][0]
liq2 = [shape for shape in slide.shapes if shape.name == "liq2"][0]
liq3 = [shape for shape in slide.shapes if shape.name == "liq3"][0]
liq4 = [shape for shape in slide.shapes if shape.name == "liq4"][0]
regua = [shape for shape in slide.shapes if shape.name == "regua"][0]
carteira_vs_bench = [
    shape for shape in slide.shapes if shape.name == "carteira_vs_bench"][0].chart
for i in range(5):
    for j in range(5):
        carteiraBenchmark_on_table.cell(
            i+2, j+1).text_frame.paragraphs[0].runs[0].text = final_onshore["df_final_CarteiraBenchmark"].iloc[i, j+1]

liq_val1.text = final_onshore["df_final_FluxoResgate"].iloc[0, 1]
liq1.text_frame.paragraphs[0].runs[0].text = f"{int(np.round(final_onshore['df_final_FluxoResgate'].iloc[0,2]))}%" if final_onshore[
    'df_final_FluxoResgate'].iloc[0, 2] > 10 else ""
liq1.width = int(
    regua.width*final_onshore['df_final_FluxoResgate'].iloc[0, 2]/100)
liq_val2.text = final_onshore["df_final_FluxoResgate"].iloc[1, 1]

liq2.text_frame.paragraphs[0].runs[0].text = f"{int(np.round(final_onshore['df_final_FluxoResgate'].iloc[1,2]))}%" if final_onshore[
    'df_final_FluxoResgate'].iloc[1, 2] > 10 else ""
liq2.left = liq1.left + liq1.width
liq2.width = int(
    regua.width*final_onshore['df_final_FluxoResgate'].iloc[1, 2]/100)
liq_val3.text = final_onshore["df_final_FluxoResgate"].iloc[2, 1]

liq3.text_frame.paragraphs[0].runs[0].text = f"{int(np.round(final_onshore['df_final_FluxoResgate'].iloc[2,2]))}%" if final_onshore[
    'df_final_FluxoResgate'].iloc[2, 2] > 10 else ""
liq3.left = liq2.left + liq2.width
liq3.width = int(
    regua.width*final_onshore['df_final_FluxoResgate'].iloc[2, 2]/100)
liq_val4.text = final_onshore["df_final_FluxoResgate"].iloc[3, 1]

liq4.text_frame.paragraphs[0].runs[0].text = f"{int(np.round(final_onshore['df_final_FluxoResgate'].iloc[3,2]))}%" if final_onshore[
    'df_final_FluxoResgate'].iloc[3, 2] > 10 else ""
liq4.left = liq3.left + liq3.width
liq4.width = int(
    regua.width*final_onshore['df_final_FluxoResgate'].iloc[3, 2]/100)

# ---define new chart data---
chart_data = CategoryChartData()
chart_data.categories = [
    date for date in final_onshore["df_final_RendimentoCarteira"]["Data"]]
chart_data.add_series('Carteira', tuple(
    final_onshore["df_final_RendimentoCarteira"]["Carteira"].ffill().fillna(0).values))
chart_data.add_series('CDI', tuple(
    final_onshore["df_final_RendimentoCarteira"]["CDI"].ffill().fillna(0).values))
# ---replace chart data---
carteira_vs_bench.replace_data(chart_data)

# slide 4 - Historico e Resumo por Conta (onshore)
slide = prs.slides[4]  # Acessa o slide desejado
table_retorno_consolidado = [
    shape for shape in slide.shapes if shape.name == "retorno_consolidado"][0].table

for row in range(len(final_onshore["df_final_RetornosNominais"])):
    for column in range(16):
        table_retorno_consolidado.cell(
            2+row, column).text = final_onshore["df_final_RetornosNominais"].iloc[row, column]
        if (row % 2):
            # bench
            cell = table_retorno_consolidado.cell(2+row, column)
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0xf8, 0xf8, 0xf8)
            if column == 0:
                cell.text = ""

            cell.text_frame.paragraphs[0].font.size = Pt(7)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            table_retorno_consolidado.rows[2+row].height = 216000

        else:
            # carteira
            cell = table_retorno_consolidado.cell(2+row, column)
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)

            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            table_retorno_consolidado.rows[2+row].height = 288000

for _ in range(row+3, 22):
    remove_row(table_retorno_consolidado,
                table_retorno_consolidado.rows[row+3])

visao_custodiante = [
    shape for shape in slide.shapes if shape.name == "visao_custodiante"][0]
table_visao_custodiante = visao_custodiante.table
total = [shape for shape in slide.shapes if shape.name == "total"][0]
table_total = total.table
total_onshore = cmd_onshore["df_ResumoPorConta"].sum()
for i in range(len(total_onshore)-3):
    table_total.cell(
        0, i+1).text_frame.paragraphs[0].runs[0].text = round_if_numeric(total_onshore[i+2], 0, "R$ ")
table_total.cell(0, 5).text_frame.paragraphs[0].runs[0].text = round_if_numeric(
    total_onshore[6], 1, "%", False)

# if len(final_onshore["df_final_ResumoPorConta"]) > 1:
#         for _ in range(2,len(final_onshore["df_final_ResumoPorConta"])):
#             add_row(table_visao_custodiante)

for row in range(len(final_onshore["df_final_ResumoPorConta"])):
    for column in range(6):
        table_visao_custodiante.cell(
            2+row, column).text = final_onshore["df_final_ResumoPorConta"].iloc[row, column+1]

        if (row % 2):
            # segundo
            cell = table_visao_custodiante.cell(2+row, column)
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0xf8, 0xf8, 0xf8)

            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            table_visao_custodiante.rows[2+row].height = 288000
        else:
            # primeiro
            cell = table_visao_custodiante.cell(2+row, column)
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)

            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            table_visao_custodiante.rows[2+row].height = 288000

        if column == 0:
            cell = table_visao_custodiante.cell(2+row, column)
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].font.bold = True

for _ in range(row+3, 9):
    remove_row(table_visao_custodiante,
                table_visao_custodiante.rows[row+3])

retorno_consolidado_height = (
    (len(table_retorno_consolidado.rows)-2)/2)*Cm(1.4) + Cm(1.88)
visao_custodiante_height = (
    len(table_visao_custodiante.rows)-2)*Cm(0.8) + Cm(1.8)
total_height = retorno_consolidado_height + \
    visao_custodiante_height + Cm(0.5)
slide_height = Cm(19.05)
[shape for shape in slide.shapes if shape.name ==
    "retorno_consolidado"][0].top = int(slide_height/2 - total_height/2)
visao_custodiante.top = int([shape for shape in slide.shapes if shape.name ==
                            "retorno_consolidado"][0].top + retorno_consolidado_height + Cm(0.5))

total.top = visao_custodiante.top + visao_custodiante_height




# slide 6 - Retorno por classe (Onshore)
slide = prs.slides[5]
alocacao_classe = [
    shape for shape in slide.shapes if shape.name == "alocacao_classe"][0].chart


rentabilidade_classe= [shape for shape in slide.shapes if shape.name == "rentabilidade_classe"][0].table

perf_attr_fill(final_onshore["df_final_AlocacaoClasse"]["Classe"].to_list(
), final_onshore["df_final_PerfAttr"].copy(), slide)

# ---define new chart data---
chart_data = CategoryChartData()
chart_data.categories = [
    date for date in final_onshore["df_final_AlocacaoClasse"]["Classe"]]
chart_data.add_series('% Carteira', tuple(
    final_onshore["df_final_AlocacaoClasse"]["% da Carteira"].values))
# ---replace chart data---
alocacao_classe.replace_data(chart_data)

for row in range(len(cmd_onshore["df_ClasseDeAtivos"])):
    for column in range(7):
        rentabilidade_classe.cell(1+row, column).text_frame.paragraphs[0].runs[0].text = round_if_numeric(cmd_onshore["df_ClasseDeAtivos"].replace("nd","-").iloc[row,column],1,"%",False)
for _ in range(len(cmd_onshore["df_ClasseDeAtivos"]),9):
    remove_row(rentabilidade_classe, rentabilidade_classe.rows[len(cmd_onshore["df_ClasseDeAtivos"])+1])


#tabela boni Onshore
slide = prs.slides[6]
table_boni_mes_onshore = [shape for shape in slide.shapes if shape.name == "boni_mes"][0].table
table_boni_ano_onshore = [shape for shape in slide.shapes if shape.name == "boni_ano"][0].table

for row in range(len(carteira_ekho_onshore_mes)):
    table_boni_mes_onshore.cell(3+row, 0).text_frame.paragraphs[0].runs[0].text = carteira_ekho_onshore_mes.iloc[row,0] #Classe
    table_boni_mes_onshore.cell(3+row, 2).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_mes.iloc[row,1],0,"%",False) #% Alocado
    table_boni_mes_onshore.cell(3+row, 3).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_mes.iloc[row,2],0,"R$ ",True) #Saldo Bruto
    table_boni_mes_onshore.cell(3+row, 4).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_mes.iloc[row,3],1,"%",False) #Performance
    table_boni_mes_onshore.cell(3+row, 5).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_mes.iloc[row,4],0,"%",False) #Performance Relativa
    table_boni_mes_onshore.cell(3+row, 6).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_mes.iloc[row,5],1,"%",False) #Contribuiçao Bruta (%)
    table_boni_mes_onshore.cell(3+row, 7).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_mes.iloc[row,6],0,"R$ ",True) #Contribuiçao Bruta (Dinheiro)
    table_boni_mes_onshore.cell(3+row, 9).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_mes.iloc[row,7],0,"%",False) #Peso PL
    table_boni_mes_onshore.cell(3+row, 10).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_mes.iloc[row,8],1,"%",False) #Performance Benchmark
    table_boni_mes_onshore.cell(3+row, 11).text_frame.paragraphs[0].runs[0].text = carteira_ekho_onshore_mes.iloc[row,9] #Benchmark
    table_boni_mes_onshore.cell(3+row, 12).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_mes.iloc[row,10],1,"%",False) #Contribuicao
    table_boni_mes_onshore.cell(3+row, 14).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_mes.iloc[row,11],1,"%",False) #Total
    table_boni_mes_onshore.cell(3+row, 15).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_mes.iloc[row,12],1,"%",False) #Alocacao
    table_boni_mes_onshore.cell(3+row, 16).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_mes.iloc[row,13],1,"%",False) #Selecao
    
table_boni_mes_onshore.cell(10, 2).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_mes.iloc[:,1].sum(),0,"%",False) # % Alocado
table_boni_mes_onshore.cell(10, 3).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_mes.iloc[:,2].sum(),0,"R$ ",True) #Saldo Bruto
table_boni_mes_onshore.cell(10, 4).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_mes.iloc[:,5].sum(),1,"%",False) #Performance
table_boni_mes_onshore.cell(10, 6).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_mes.iloc[:,5].sum(),1,"%",False) #Contribuiçao Bruta (%)
table_boni_mes_onshore.cell(10, 7).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_mes.iloc[:,6].sum(),0,"R$ ",True) #Contribuiçao Bruta (Dinheiro)
table_boni_mes_onshore.cell(10, 9).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_mes.iloc[:,7].sum(),0,"%",False) #Peso PL
table_boni_mes_onshore.cell(10, 10).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_mes.iloc[:,10].sum(),1,"%",False) #Performance Benchmark
table_boni_mes_onshore.cell(10, 12).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_mes.iloc[:,10].sum(),1,"%",False) #Contribuicao
table_boni_mes_onshore.cell(10, 14).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_mes.iloc[:,11].sum(),1,"%",False) #Total
table_boni_mes_onshore.cell(10, 15).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_mes.iloc[:,12].sum(),1,"%",False) #Alocacao
table_boni_mes_onshore.cell(10, 16).text_frame.paragraphs[0].runs[0].text = round_if_numeric((carteira_ekho_onshore_mes.iloc[:,11].sum())-(carteira_ekho_onshore_mes.iloc[:,12].sum()),1,"%",False) #Selecao
table_boni_mes_onshore.cell(10, 5).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_mes.iloc[:,10].sum()/carteira_ekho_onshore_mes.iloc[:,5].sum(),0,"%",False) #Performance Relativa


for row in range(len(carteira_ekho_onshore_ano)):
    table_boni_ano_onshore.cell(3+row, 0).text_frame.paragraphs[0].runs[0].text = carteira_ekho_onshore_ano.iloc[row,0] #Classe
    table_boni_ano_onshore.cell(3+row, 2).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_ano.iloc[row,1],0,"%",False) #% Alocado
    table_boni_ano_onshore.cell(3+row, 3).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_ano.iloc[row,2],0,"R$ ",True) #Saldo Bruto
    table_boni_ano_onshore.cell(3+row, 4).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_ano.iloc[row,3],1,"%",False) #Performance
    table_boni_ano_onshore.cell(3+row, 5).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_ano.iloc[row,4],0,"%",False) #Performance Relativa
    table_boni_ano_onshore.cell(3+row, 6).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_ano.iloc[row,5],1,"%",False) #Contribuiçao Bruta (%)
    table_boni_ano_onshore.cell(3+row, 7).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_ano.iloc[row,6],0,"R$ ",True) #Contribuiçao Bruta (Dinheiro)
    table_boni_ano_onshore.cell(3+row, 9).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_ano.iloc[row,7],0,"%",False) #Peso PL
    table_boni_ano_onshore.cell(3+row, 10).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_ano.iloc[row,8],1,"%",False) #Performance Benchmark
    table_boni_ano_onshore.cell(3+row, 11).text_frame.paragraphs[0].runs[0].text = carteira_ekho_onshore_ano.iloc[row,9] #Benchmark
    table_boni_ano_onshore.cell(3+row, 12).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_ano.iloc[row,10],1,"%",False) #Contribuicao
    table_boni_ano_onshore.cell(3+row, 14).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_ano.iloc[row,11],1,"%",False) #Total
    table_boni_ano_onshore.cell(3+row, 15).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_ano.iloc[row,12],1,"%",False) #Alocacao
    table_boni_ano_onshore.cell(3+row, 16).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_ano.iloc[row,13],1,"%",False) #Selecao
    
table_boni_ano_onshore.cell(10, 2).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_ano.iloc[:,1].sum(),0,"%",False) # % Alocado
table_boni_ano_onshore.cell(10, 3).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_ano.iloc[:,2].sum(),0,"R$ ",True) #Saldo Bruto
table_boni_ano_onshore.cell(10, 4).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_ano.iloc[:,5].sum(),1,"%",False) #Performance
table_boni_ano_onshore.cell(10, 6).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_ano.iloc[:,5].sum(),1,"%",False) #Contribuiçao Bruta (%)
table_boni_ano_onshore.cell(10, 7).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_ano.iloc[:,6].sum(),0,"R$ ",True) #Contribuiçao Bruta (Dinheiro)
table_boni_ano_onshore.cell(10, 9).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_ano.iloc[:,7].sum(),0,"%",False) #Peso PL
table_boni_ano_onshore.cell(10, 10).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_ano.iloc[:,10].sum(),1,"%",False) #Performance Benchmark
table_boni_ano_onshore.cell(10, 12).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_ano.iloc[:,10].sum(),1,"%",False) #Contribuicao
table_boni_ano_onshore.cell(10, 14).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_ano.iloc[:,11].sum(),1,"%",False) #Total
table_boni_ano_onshore.cell(10, 15).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_ano.iloc[:,12].sum(),1,"%",False) #Alocacao
table_boni_ano_onshore.cell(10, 16).text_frame.paragraphs[0].runs[0].text = round_if_numeric((carteira_ekho_onshore_ano.iloc[:,11].sum())-(carteira_ekho_onshore_ano.iloc[:,12].sum()),1,"%",False) #Selecao
table_boni_ano_onshore.cell(10, 5).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_onshore_ano.iloc[:,10].sum()/carteira_ekho_onshore_ano.iloc[:,5].sum(),0,"%",False) #Performance Relativa

# tabela boni Offshore
slide = prs.slides[25]
table_boni_mes_offshore = [shape for shape in slide.shapes if shape.name == "boni_mes"][0].table
table_boni_ano_offshore = [shape for shape in slide.shapes if shape.name == "boni_ano"][0].table

for row in range(len(carteira_ekho_offshore_mes)):
    table_boni_mes_offshore.cell(3+row, 0).text_frame.paragraphs[0].runs[0].text = carteira_ekho_offshore_mes.iloc[row,0] #Classe
    table_boni_mes_offshore.cell(3+row, 2).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_mes.iloc[row,1],0,"%",False) #% Alocado
    table_boni_mes_offshore.cell(3+row, 3).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_mes.iloc[row,2],0,"$ ",True) #Saldo Bruto
    table_boni_mes_offshore.cell(3+row, 4).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_mes.iloc[row,3],1,"%",False) #Performance
    table_boni_mes_offshore.cell(3+row, 5).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_mes.iloc[row,4],0,"%",False) #Performance Relativa
    table_boni_mes_offshore.cell(3+row, 6).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_mes.iloc[row,5],1,"%",False) #Contribuiçao Bruta (%)
    table_boni_mes_offshore.cell(3+row, 7).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_mes.iloc[row,6],0,"$ ",True) #Contribuiçao Bruta (Dinheiro)
    table_boni_mes_offshore.cell(3+row, 9).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_mes.iloc[row,7],0,"%",False) #Peso PL
    table_boni_mes_offshore.cell(3+row, 10).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_mes.iloc[row,8],1,"%",False) #Performance Benchmark
    table_boni_mes_offshore.cell(3+row, 11).text_frame.paragraphs[0].runs[0].text = carteira_ekho_offshore_mes.iloc[row,9] #Benchmark
    table_boni_mes_offshore.cell(3+row, 12).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_mes.iloc[row,10],1,"%",False) #Contribuicao
    table_boni_mes_offshore.cell(3+row, 14).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_mes.iloc[row,11],1,"%",False) #Total
    table_boni_mes_offshore.cell(3+row, 15).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_mes.iloc[row,12],1,"%",False) #Alocacao
    table_boni_mes_offshore.cell(3+row, 16).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_mes.iloc[row,13],1,"%",False) #Selecao
    
table_boni_mes_offshore.cell(12, 2).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_mes.iloc[:,1].sum(),0,"%",False) # % Alocado
table_boni_mes_offshore.cell(12, 3).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_mes.iloc[:,2].sum(),0,"$ ",True) #Saldo Bruto
table_boni_mes_offshore.cell(12, 4).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_mes.iloc[:,5].sum(),1,"%",False) #Performance
table_boni_mes_offshore.cell(12, 6).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_mes.iloc[:,5].sum(),1,"%",False) #Contribuiçao Bruta (%)
table_boni_mes_offshore.cell(12, 7).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_mes.iloc[:,6].sum(),0,"$ ",True) #Contribuiçao Bruta (Dinheiro)
table_boni_mes_offshore.cell(12, 9).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_mes.iloc[:,7].sum(),0,"%",False) #Peso PL
table_boni_mes_offshore.cell(12, 10).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_mes.iloc[:,10].sum(),1,"%",False) #Performance Benchmark
table_boni_mes_offshore.cell(12, 12).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_mes.iloc[:,10].sum(),1,"%",False) #Contribuicao
table_boni_mes_offshore.cell(12, 14).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_mes.iloc[:,11].sum(),1,"%",False) #Total
table_boni_mes_offshore.cell(12, 15).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_mes.iloc[:,12].sum(),1,"%",False) #Alocacao
table_boni_mes_offshore.cell(12, 16).text_frame.paragraphs[0].runs[0].text = round_if_numeric((carteira_ekho_offshore_mes.iloc[:,11].sum())-(carteira_ekho_offshore_mes.iloc[:,12].sum()),1,"%",False) #Selecao
table_boni_mes_offshore.cell(12, 5).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_mes.iloc[:,10].sum()/carteira_ekho_offshore_mes.iloc[:,5].sum(),0,"%",False) #Performance Relativa


for row in range(len(carteira_ekho_offshore_ano)):
    table_boni_ano_offshore.cell(3+row, 0).text_frame.paragraphs[0].runs[0].text = carteira_ekho_offshore_ano.iloc[row,0] #Classe
    table_boni_ano_offshore.cell(3+row, 2).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_ano.iloc[row,1],0,"%",False) #% Alocado
    table_boni_ano_offshore.cell(3+row, 3).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_ano.iloc[row,2],0,"$ ",True) #Saldo Bruto
    table_boni_ano_offshore.cell(3+row, 4).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_ano.iloc[row,3],1,"%",False) #Performance
    table_boni_ano_offshore.cell(3+row, 5).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_ano.iloc[row,4],0,"%",False) #Performance Relativa
    table_boni_ano_offshore.cell(3+row, 6).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_ano.iloc[row,5],1,"%",False) #Contribuiçao Bruta (%)
    table_boni_ano_offshore.cell(3+row, 7).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_ano.iloc[row,6],0,"$ ",True) #Contribuiçao Bruta (Dinheiro)
    table_boni_ano_offshore.cell(3+row, 9).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_ano.iloc[row,7],0,"%",False) #Peso PL
    table_boni_ano_offshore.cell(3+row, 10).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_ano.iloc[row,8],1,"%",False) #Performance Benchmark
    table_boni_ano_offshore.cell(3+row, 11).text_frame.paragraphs[0].runs[0].text = carteira_ekho_offshore_ano.iloc[row,9] #Benchmark
    table_boni_ano_offshore.cell(3+row, 12).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_ano.iloc[row,10],1,"%",False) #Contribuicao
    table_boni_ano_offshore.cell(3+row, 14).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_ano.iloc[row,11],1,"%",False) #Total
    table_boni_ano_offshore.cell(3+row, 15).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_ano.iloc[row,12],1,"%",False) #Alocacao
    table_boni_ano_offshore.cell(3+row, 16).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_ano.iloc[row,13],1,"%",False) #Selecao
    
table_boni_ano_offshore.cell(12, 2).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_ano.iloc[:,1].sum(),0,"%",False) # % Alocado
table_boni_ano_offshore.cell(12, 3).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_ano.iloc[:,2].sum(),0,"$ ",True) #Saldo Bruto
table_boni_ano_offshore.cell(12, 4).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_ano.iloc[:,5].sum(),1,"%",False) #Performance
table_boni_ano_offshore.cell(12, 6).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_ano.iloc[:,5].sum(),1,"%",False) #Contribuiçao Bruta (%)
table_boni_ano_offshore.cell(12, 7).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_ano.iloc[:,6].sum(),0,"$ ",True) #Contribuiçao Bruta (Dinheiro)
table_boni_ano_offshore.cell(12, 9).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_ano.iloc[:,7].sum(),0,"%",False) #Peso PL
table_boni_ano_offshore.cell(12, 10).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_ano.iloc[:,10].sum(),1,"%",False) #Performance Benchmark
table_boni_ano_offshore.cell(12, 12).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_ano.iloc[:,10].sum(),1,"%",False) #Contribuicao
table_boni_ano_offshore.cell(12, 14).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_ano.iloc[:,11].sum(),1,"%",False) #Total
table_boni_ano_offshore.cell(12, 15).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_ano.iloc[:,12].sum(),1,"%",False) #Alocacao
table_boni_ano_offshore.cell(12, 16).text_frame.paragraphs[0].runs[0].text = round_if_numeric((carteira_ekho_offshore_ano.iloc[:,11].sum())-(carteira_ekho_offshore_ano.iloc[:,12].sum()),1,"%",False) #Selecao
table_boni_ano_offshore.cell(12, 5).text_frame.paragraphs[0].runs[0].text = round_if_numeric(carteira_ekho_offshore_ano.iloc[:,10].sum()/carteira_ekho_offshore_ano.iloc[:,5].sum(),0,"%",False) #Performance Relativa


# slide 17 - Historico e Resumo por Conta (offshore)

slide = prs.slides[23]  # Acessa o slide desejado
table_retorno_consolidado = [
    shape for shape in slide.shapes if shape.name == "retorno_consolidado"][0].table

for row in range(len(final_offshore["df_final_RetornosNominais"])):
    for column in range(16):
        table_retorno_consolidado.cell(
            2+row, column).text = final_offshore["df_final_RetornosNominais"].iloc[row, column]
        if (row % 2):
            # bench
            cell = table_retorno_consolidado.cell(2+row, column)
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0xf8, 0xf8, 0xf8)
            if column == 0:
                cell.text = ""

            cell.text_frame.paragraphs[0].font.size = Pt(7)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            table_retorno_consolidado.rows[2+row].height = 216000

        else:
            # carteira
            cell = table_retorno_consolidado.cell(2+row, column)
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)

            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            table_retorno_consolidado.rows[2+row].height = 288000

for _ in range(row+3, 22):
    remove_row(table_retorno_consolidado,
                table_retorno_consolidado.rows[row+3])

visao_custodiante = [
    shape for shape in slide.shapes if shape.name == "visao_custodiante"][0]
table_visao_custodiante = visao_custodiante.table
total = [shape for shape in slide.shapes if shape.name == "total"][0]
table_total = total.table
total_offshore = cmd_offshore["df_ResumoPorConta"].sum()
for i in range(len(total_offshore)-3):
    table_total.cell(
        0, i+1).text_frame.paragraphs[0].runs[0].text = round_if_numeric(total_offshore[i+2], 0, "$ ")
table_total.cell(0, 5).text_frame.paragraphs[0].runs[0].text = round_if_numeric(
    total_offshore[6], 1, "%", False)

# if len(final_offshore["df_final_ResumoPorConta"]) > 1:
#         for _ in range(2,len(final_offshore["df_final_ResumoPorConta"])):
#             add_row(table_visao_custodiante)

for row in range(len(final_offshore["df_final_ResumoPorConta"])):
    for column in range(6):
        table_visao_custodiante.cell(
            2+row, column).text = final_offshore["df_final_ResumoPorConta"].iloc[row, column+1]

        if (row % 2):
            # segundo
            cell = table_visao_custodiante.cell(2+row, column)
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0xf8, 0xf8, 0xf8)

            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            table_visao_custodiante.rows[2+row].height = 288000
        else:
            # primeiro
            cell = table_visao_custodiante.cell(2+row, column)
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)

            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            table_visao_custodiante.rows[2+row].height = 288000

        if column == 0:
            cell = table_visao_custodiante.cell(2+row, column)
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].font.bold = True

for _ in range(row+3, 9):
    remove_row(table_visao_custodiante,
                table_visao_custodiante.rows[row+3])

retorno_consolidado_height = (
    (len(table_retorno_consolidado.rows)-2)/2)*Cm(1.4) + Cm(1.88)
visao_custodiante_height = (
    len(table_visao_custodiante.rows)-2)*Cm(0.8) + Cm(1.8)
total_height = retorno_consolidado_height + \
    visao_custodiante_height + Cm(0.5)
slide_height = Cm(19.05)
[shape for shape in slide.shapes if shape.name ==
    "retorno_consolidado"][0].top = int(slide_height/2 - total_height/2)
visao_custodiante.top = int([shape for shape in slide.shapes if shape.name ==
                            "retorno_consolidado"][0].top + retorno_consolidado_height + Cm(0.5))

total.top = visao_custodiante.top + visao_custodiante_height

# slide 18 - Retorno por classe (Offshore)
slide = prs.slides[24]
alocacao_classe = [
    shape for shape in slide.shapes if shape.name == "alocacao_classe"][0].chart


rentabilidade_classe= [shape for shape in slide.shapes if shape.name == "rentabilidade_classe"][0].table

perf_attr_fill(final_offshore["df_final_AlocacaoClasse"]["Classe"].to_list(
), final_offshore["df_final_PerfAttr"].copy(), slide)

# ---define new chart data---
chart_data = CategoryChartData()
chart_data.categories = [
    date for date in final_offshore["df_final_AlocacaoClasse"]["Classe"]]
chart_data.add_series('% Carteira', tuple(
    final_offshore["df_final_AlocacaoClasse"]["% da Carteira"].values))
# ---replace chart data---
alocacao_classe.replace_data(chart_data)

for row in range(len(cmd_offshore["df_ClasseDeAtivos"])):
    for column in range(7):
        rentabilidade_classe.cell(1+row, column).text_frame.paragraphs[0].runs[0].text = round_if_numeric(cmd_offshore["df_ClasseDeAtivos"].replace("nd","-").iloc[row,column],1,"%",False)
for _ in range(len(cmd_offshore["df_ClasseDeAtivos"]),9):
    remove_row(rentabilidade_classe, rentabilidade_classe.rows[len(cmd_offshore["df_ClasseDeAtivos"])+1])


# slide 16 - Resumo Offshore
slide = prs.slides[22]
carteiraBenchmark_on_table = [
    shape for shape in slide.shapes if shape.name == "CarteiraBenchmark"][0].table
liq_val1 = [shape for shape in slide.shapes if shape.name ==
            "liq_val1"][0].text_frame.paragraphs[0].runs[0]
liq_val2 = [shape for shape in slide.shapes if shape.name ==
            "liq_val2"][0].text_frame.paragraphs[0].runs[0]
liq_val3 = [shape for shape in slide.shapes if shape.name ==
            "liq_val3"][0].text_frame.paragraphs[0].runs[0]
liq_val4 = [shape for shape in slide.shapes if shape.name ==
            "liq_val4"][0].text_frame.paragraphs[0].runs[0]
liq1 = [shape for shape in slide.shapes if shape.name == "liq1"][0]
liq2 = [shape for shape in slide.shapes if shape.name == "liq2"][0]
liq3 = [shape for shape in slide.shapes if shape.name == "liq3"][0]
liq4 = [shape for shape in slide.shapes if shape.name == "liq4"][0]
regua = [shape for shape in slide.shapes if shape.name == "regua"][0]
carteira_vs_bench = [
    shape for shape in slide.shapes if shape.name == "carteira_vs_bench"][0].chart
for i in range(5):
    for j in range(5):
        carteiraBenchmark_on_table.cell(
            i+2, j+1).text_frame.paragraphs[0].runs[0].text = final_offshore["df_final_CarteiraBenchmark"].iloc[i, j+1]

liq_val1.text = final_offshore["df_final_FluxoResgate"].iloc[0, 1]
liq1.text_frame.paragraphs[0].runs[0].text = f"{int(np.round(final_offshore['df_final_FluxoResgate'].iloc[0,2]))}%" if final_offshore[
    'df_final_FluxoResgate'].iloc[0, 2] > 10 else ""
liq1.width = int(
    regua.width*final_offshore['df_final_FluxoResgate'].iloc[0, 2]/100)
liq_val2.text = final_offshore["df_final_FluxoResgate"].iloc[1, 1]

liq2.text_frame.paragraphs[0].runs[0].text = f"{int(np.round(final_offshore['df_final_FluxoResgate'].iloc[1,2]))}%" if final_offshore[
    'df_final_FluxoResgate'].iloc[1, 2] > 10 else ""
liq2.left = liq1.left + liq1.width
liq2.width = int(
    regua.width*final_offshore['df_final_FluxoResgate'].iloc[1, 2]/100)
liq_val3.text = final_offshore["df_final_FluxoResgate"].iloc[2, 1]

liq3.text_frame.paragraphs[0].runs[0].text = f"{int(np.round(final_offshore['df_final_FluxoResgate'].iloc[2,2]))}%" if final_offshore[
    'df_final_FluxoResgate'].iloc[2, 2] > 10 else ""
liq3.left = liq2.left + liq2.width
liq3.width = int(
    regua.width*final_offshore['df_final_FluxoResgate'].iloc[2, 2]/100)
liq_val4.text = final_offshore["df_final_FluxoResgate"].iloc[3, 1]

liq4.text_frame.paragraphs[0].runs[0].text = f"{int(np.round(final_offshore['df_final_FluxoResgate'].iloc[3,2]))}%" if final_offshore[
    'df_final_FluxoResgate'].iloc[3, 2] > 10 else ""
liq4.left = liq3.left + liq3.width
liq4.width = int(
    regua.width*final_offshore['df_final_FluxoResgate'].iloc[3, 2]/100)

# ---define new chart data---
chart_data = CategoryChartData()
chart_data.categories = [
    date for date in final_offshore["df_final_RendimentoCarteira"]["Data"]]
chart_data.add_series('Carteira', tuple(
    final_offshore["df_final_RendimentoCarteira"]["Carteira"].ffill().fillna(0).values))
chart_data.add_series('EFFR', tuple(
    final_offshore["df_final_RendimentoCarteira"]["EFFR"].ffill().fillna(0).values))
# ---replace chart data---
carteira_vs_bench.replace_data(chart_data)

# Fill onshore tables
max_slide_onshore = fill_asset_table(
    7, final_onshore["dfs_class"], 18)

# Fill offshore tables
max_slide_offshore = fill_asset_table(
    26, final_offshore["dfs_class"], 18)

# Delete unused slides
for slide_index in range(max_slide_offshore, 39):
    delete_slide(prs, max_slide_offshore+1)

for slide_index in range(max_slide_onshore, 20):
    delete_slide(prs, max_slide_onshore+1)
    

# slide_number
for slide_idx in range(len(prs.slides)):
    slide = prs.slides[slide_idx]
    try:
        name_shape = [
            shape for shape in slide.shapes if shape.name == "slide_number"][0]
        name_shape.text_frame.paragraphs[0].runs[0].text = str(slide_idx+1)
    except:
        pass

prs.save(
    f"{folder}/{client_name}_{meses[data_fim.month]}2024_aaaa.pptx")
print(f"{client_name} report DONE.")